#!/usr/bin/env python3
from __future__ import annotations

import argparse
import re
from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR_TYPE
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


REPO_ROOT = Path(__file__).resolve().parents[1]
SOURCE_MD = REPO_ROOT / "presentations" / "prompt-to-skill-talk.md"
DEFAULT_OUTPUT = REPO_ROOT / "presentations" / "prompt-to-skill-talk-boeet-generated.pptx"
DEFAULT_TEMPLATE = Path.home() / ".agents" / "skills" / "boeet-pptx" / "assets" / "boeet-template-16x9-internal.pptx"

FONT = "微软雅黑"
MONO = "微软雅黑"

BLUE = RGBColor(0x00, 0x58, 0x89)
SKY = RGBColor(0x00, 0x9A, 0xD8)
GREEN = RGBColor(0x45, 0xBA, 0x7D)
CORAL = RGBColor(0xFF, 0x8F, 0x82)
YELLOW = RGBColor(0xFF, 0xC7, 0x57)
DARK = RGBColor(0x54, 0x58, 0x5A)
DEEP = RGBColor(0x37, 0x35, 0x45)
INK = RGBColor(0x1F, 0x2D, 0x38)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF4, 0xF8, 0xFB)
MIST = RGBColor(0xE7, 0xF1, 0xF7)
PALE_GREEN = RGBColor(0xE7, 0xF6, 0xEE)
PALE_YELLOW = RGBColor(0xFF, 0xF6, 0xDB)
PALE_CORAL = RGBColor(0xFF, 0xEE, 0xEA)
SOFT_LINE = RGBColor(0xC8, 0xD9, 0xE6)


@dataclass
class Section:
    number: int
    title: str
    body: str


def parse_sections(markdown: str) -> dict[int, Section]:
    pattern = re.compile(r"^### P(\d+)｜(.+?)\n(.*?)(?=^### P\d+｜|\Z)", re.M | re.S)
    sections: dict[int, Section] = {}
    for match in pattern.finditer(markdown):
        number = int(match.group(1))
        sections[number] = Section(number=number, title=match.group(2).strip(), body=match.group(3).strip())
    expected = list(range(1, 17))
    missing = [number for number in expected if number not in sections]
    if missing:
        raise ValueError(f"Missing slide sections in source markdown: {missing}")
    return sections


def set_run_style(run, size: int, *, color: RGBColor = DARK, bold: bool = False, italic: bool = False, font_name: str = FONT) -> None:
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic


def set_paragraph(
    paragraph,
    text: str,
    *,
    size: int,
    color: RGBColor = DARK,
    bold: bool = False,
    italic: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    font_name: str = FONT,
) -> None:
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    set_run_style(run, size=size, color=color, bold=bold, italic=italic, font_name=font_name)


def clear_text_shape(shape) -> None:
    if not shape.has_text_frame:
        return
    shape.text_frame.clear()


def get_placeholder(slide, idx: int):
    for shape in slide.shapes:
        placeholder = getattr(shape, "placeholder_format", None)
        if placeholder and placeholder.idx == idx:
            return shape
    raise KeyError(f"Placeholder {idx} not found on slide.")


def set_placeholder_text(slide, idx: int, text: str, *, size: int = 12, color: RGBColor = DARK, bold: bool = False, align: PP_ALIGN = PP_ALIGN.LEFT) -> None:
    shape = get_placeholder(slide, idx)
    clear_text_shape(shape)
    frame = shape.text_frame
    frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    set_paragraph(paragraph, text, size=size, color=color, bold=bold, align=align)


def add_textbox(
    slide,
    x,
    y,
    w,
    h,
    text: str,
    *,
    size: int = 18,
    color: RGBColor = DARK,
    bold: bool = False,
    italic: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    font_name: str = FONT,
    fill_color: RGBColor | None = None,
    line_color: RGBColor | None = None,
    radius: str = "none",
    margin: float = 0.08,
):
    if radius == "rounded":
        box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
        box.adjustments[0] = 0.18
    else:
        box = slide.shapes.add_textbox(x, y, w, h)
    if fill_color is not None:
        box.fill.solid()
        box.fill.fore_color.rgb = fill_color
    elif hasattr(box, "fill"):
        box.fill.background()
    if line_color is not None:
        box.line.color.rgb = line_color
        box.line.width = Pt(1.2)
    elif hasattr(box, "line"):
        box.line.fill.background()
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    paragraph = frame.paragraphs[0]
    set_paragraph(paragraph, text, size=size, color=color, bold=bold, italic=italic, align=align, font_name=font_name)
    return box


def add_lines(
    slide,
    x,
    y,
    w,
    h,
    lines: list[str],
    *,
    size: int = 13,
    color: RGBColor = DARK,
    bold_first: bool = False,
    spacing: int = 6,
    prefix: str = "",
    font_name: str = FONT,
):
    box = slide.shapes.add_textbox(x, y, w, h)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.02)
    frame.margin_right = Inches(0.02)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    for index, line in enumerate(lines):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        text = f"{prefix}{line}" if prefix else line
        set_paragraph(paragraph, text, size=size, color=color, bold=bold_first and index == 0, font_name=font_name)
        paragraph.space_after = Pt(spacing)
    return box


def add_card(
    slide,
    x,
    y,
    w,
    h,
    title: str,
    lines: list[str],
    *,
    accent: RGBColor,
    fill: RGBColor = WHITE,
    title_size: int = 14,
    body_size: int = 11,
    tag_text: str | None = None,
):
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.adjustments[0] = 0.16
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    card.line.color.rgb = accent
    card.line.width = Pt(1.5)
    if tag_text:
        tag = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x + Inches(0.16), y + Inches(0.12), Inches(0.9), Inches(0.28))
        tag.adjustments[0] = 0.22
        tag.fill.solid()
        tag.fill.fore_color.rgb = accent
        tag.line.fill.background()
        add_textbox(slide, x + Inches(0.19), y + Inches(0.12), Inches(0.84), Inches(0.26), tag_text, size=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        title_y = y + Inches(0.48)
    else:
        title_y = y + Inches(0.18)
    add_textbox(slide, x + Inches(0.16), title_y, w - Inches(0.32), Inches(0.34), title, size=title_size, color=accent, bold=True)
    add_lines(slide, x + Inches(0.16), title_y + Inches(0.38), w - Inches(0.32), h - Inches(0.54), lines, size=body_size, color=INK, spacing=4)
    return card


def add_band(slide, x, y, w, h, text: str, *, fill: RGBColor, line: RGBColor, color: RGBColor = DARK, size: int = 12, bold: bool = False):
    return add_textbox(slide, x, y, w, h, text, size=size, color=color, bold=bold, align=PP_ALIGN.CENTER, fill_color=fill, line_color=line, radius="rounded", margin=0.06)


def add_circle(slide, x, y, diameter, text: str, *, fill: RGBColor, line: RGBColor, color: RGBColor = WHITE, size: int = 18):
    circle = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, x, y, diameter, diameter)
    circle.fill.solid()
    circle.fill.fore_color.rgb = fill
    circle.line.color.rgb = line
    circle.line.width = Pt(2)
    add_textbox(slide, x + Inches(0.08), y + Inches(0.28), diameter - Inches(0.16), diameter - Inches(0.56), text, size=size, color=color, bold=True, align=PP_ALIGN.CENTER)
    return circle


def add_connector(slide, x1, y1, x2, y2, *, color: RGBColor = SOFT_LINE, width: float = 1.4):
    connector = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, x1, y1, x2, y2)
    connector.line.color.rgb = color
    connector.line.width = Pt(width)
    return connector


def add_chevron(slide, x, y, w, h, *, fill: RGBColor = MIST):
    arrow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, x, y, w, h)
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = fill
    arrow.line.fill.background()
    return arrow


def remove_all_slides(prs: Presentation) -> None:
    slide_ids = list(prs.slides._sldIdLst)
    for slide_id in slide_ids:
        prs.part.drop_rel(slide_id.rId)
        prs.slides._sldIdLst.remove(slide_id)


def clear_content_slide(slide) -> None:
    for idx in (11, 13):
        try:
            clear_text_shape(get_placeholder(slide, idx))
        except KeyError:
            continue


def set_content_header(slide, title: str, summary: str) -> None:
    clear_content_slide(slide)
    set_placeholder_text(slide, 11, title, size=20, color=DARK, bold=True)
    set_placeholder_text(slide, 13, "")
    add_textbox(slide, Inches(0.46), Inches(0.56), Inches(11.0), Inches(0.28), summary, size=11, color=SKY)


def add_free_header(slide, title: str, summary: str) -> None:
    add_textbox(slide, Inches(0.42), Inches(0.18), Inches(10.9), Inches(0.34), title, size=21, color=DARK, bold=True)
    add_textbox(slide, Inches(0.42), Inches(0.56), Inches(10.8), Inches(0.24), summary, size=11, color=SKY)


def build_cover(slide, section: Section) -> None:
    for idx in (10, 13, 23):
        set_placeholder_text(slide, idx, "", size=12)

    add_textbox(slide, Inches(0.72), Inches(2.0), Inches(5.4), Inches(1.32), "从 Prompt 到 Skill\n把 AI 使用习惯升级为团队能力", size=24, color=DEEP, bold=True)
    add_textbox(slide, Inches(0.75), Inches(4.78), Inches(4.5), Inches(0.52), "面向混合团队的 Skill 方法、案例与落地路径", size=16, color=BLUE)
    add_textbox(slide, Inches(0.75), Inches(5.28), Inches(4.1), Inches(0.28), "研发 / 产品 / 设计 / 管理者", size=13, color=DARK)

    steps = [("Prompt", BLUE), ("Skill", GREEN), ("Team Capability", YELLOW)]
    x = Inches(6.55)
    y = Inches(2.15)
    for index, (label, color) in enumerate(steps):
        box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x + Inches(index * 1.45), y, Inches(1.22), Inches(0.58))
        box.adjustments[0] = 0.18
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = color
        box.line.width = Pt(1.6)
        add_textbox(slide, x + Inches(index * 1.45) + Inches(0.08), y + Inches(0.12), Inches(1.06), Inches(0.24), label, size=11, color=color, bold=True, align=PP_ALIGN.CENTER)
        if index < len(steps) - 1:
            add_chevron(slide, x + Inches(index * 1.45) + Inches(1.15), y + Inches(0.16), Inches(0.28), Inches(0.26), fill=MIST)
    add_band(
        slide,
        Inches(6.42),
        Inches(3.2),
        Inches(4.35),
        Inches(0.58),
        "真正值得做的，不是“更多 prompt”，而是“更稳定的工作流入口”",
        fill=LIGHT_BG,
        line=SKY,
        color=INK,
        size=11,
        bold=True,
    )
    add_textbox(slide, Inches(6.55), Inches(4.25), Inches(3.8), Inches(0.26), f"P{section.number} 封面页", size=10, color=SKY, bold=True)


def build_p2(slide, section: Section) -> None:
    set_content_header(slide, section.title, "团队已经会用 AI，但高质量做法仍停留在个人经验、聊天记录和人工兜底里。")
    cards = [
        ("反复解释", ["同一种流程每次都要重新交代", "高质量做法难以复用"], BLUE),
        ("做法漂移", ["同一种任务每个人路径不同", "质量与结构波动很大"], GREEN),
        ("经验孤岛", ["经验只留在少数人脑子里", "聊天记录无法成为团队资产"], YELLOW),
        ("风险靠盯", ["该自动守住的动作仍靠人工提醒", "一忙就容易漏"], CORAL),
    ]
    positions = [(Inches(0.5), Inches(1.35)), (Inches(3.15), Inches(1.35)), (Inches(0.5), Inches(3.15)), (Inches(3.15), Inches(3.15))]
    for (title, lines, accent), (x, y) in zip(cards, positions):
        add_card(slide, x, y, Inches(2.35), Inches(1.55), title, lines, accent=accent, fill=LIGHT_BG if accent != CORAL else PALE_CORAL)

    chain_x = Inches(6.2)
    chain_y = Inches(1.52)
    chain = [
        ("临时 prompt", BLUE),
        ("临场解释", SKY),
        ("结果漂移", GREEN),
        ("人工兜底", CORAL),
    ]
    for index, (label, color) in enumerate(chain):
        add_card(slide, chain_x, chain_y + Inches(index * 1.0), Inches(4.6), Inches(0.72), label, [], accent=color, fill=WHITE, title_size=14)
        if index < len(chain) - 1:
            add_chevron(slide, chain_x + Inches(2.08), chain_y + Inches(index * 1.0) + Inches(0.77), Inches(0.42), Inches(0.18), fill=MIST)
    add_band(slide, Inches(6.2), Inches(5.58), Inches(4.65), Inches(0.58), "症结不是“不会用 AI”，而是没把高质量做法升级成稳定能力。", fill=PALE_YELLOW, line=YELLOW, color=INK, bold=True)


def build_p3(slide, section: Section) -> None:
    set_content_header(slide, section.title, "Skill 的关键不是字数更多，而是把触发条件、步骤顺序和输出形状都固定下来。")
    add_card(
        slide,
        Inches(0.72),
        Inches(1.48),
        Inches(4.6),
        Inches(3.55),
        "Before｜一次性 prompt",
        ["每次会话都从头解释", "只对当前这一次有效", "只有熟练写 prompt 的人能稳定产出"],
        accent=BLUE,
        fill=LIGHT_BG,
        title_size=18,
        body_size=13,
    )
    add_card(
        slide,
        Inches(7.08),
        Inches(1.48),
        Inches(4.25),
        Inches(3.55),
        "After｜可复用 Skill",
        ["固定触发条件", "固定步骤顺序", "固定输出格式", "任何人都能通过统一入口复用"],
        accent=GREEN,
        fill=PALE_GREEN,
        title_size=18,
        body_size=13,
    )
    add_chevron(slide, Inches(5.62), Inches(2.7), Inches(0.82), Inches(0.42), fill=MIST)
    add_band(slide, Inches(1.2), Inches(5.35), Inches(9.8), Inches(0.64), "一句话定义：Skill 是带说明的主动工作流入口。", fill=WHITE, line=SKY, color=DEEP, size=14, bold=True)


def build_p4(slide, section: Section) -> None:
    set_content_header(slide, section.title, "Prompt 决定这次要什么，Skill 决定怎样稳定做，Rules 决定什么绝不能错。")
    columns = [
        ("Prompt", ["围绕当前请求临时组织表达", "目标是把这次需求说清楚", "更偏一次性上下文"], BLUE),
        ("Skill", ["把一类主动任务封装成流程", "强调触发条件、步骤和输出", "适合高频、可复用任务"], GREEN),
        ("Rules", ["定义项目边界、优先级和禁区", "防止越界、漏约束、误操作", "负责守住不能错的地方"], YELLOW),
    ]
    for index, (title, lines, color) in enumerate(columns):
        add_card(slide, Inches(0.65 + index * 3.72), Inches(1.55), Inches(3.15), Inches(3.8), title, lines, accent=color, fill=WHITE, title_size=18, body_size=12)
    add_band(slide, Inches(1.2), Inches(5.55), Inches(9.7), Inches(0.5), "Skill 不替代用户目标，也不替代项目规则。", fill=LIGHT_BG, line=SKY, color=INK, bold=True)


def build_p5(slide, section: Section) -> None:
    add_free_header(slide, section.title, "很多问题不必一上来就做 Agent；先看它更像边界、主动流程、自动守门，还是长链路任务。")
    add_connector(slide, Inches(1.4), Inches(5.6), Inches(10.8), Inches(5.6), color=SOFT_LINE, width=1.8)
    add_connector(slide, Inches(1.55), Inches(5.8), Inches(1.55), Inches(1.25), color=SOFT_LINE, width=1.8)
    add_chevron(slide, Inches(10.62), Inches(5.47), Inches(0.28), Inches(0.22), fill=SOFT_LINE)
    add_chevron(slide, Inches(1.42), Inches(1.16), Inches(0.22), Inches(0.28), fill=SOFT_LINE)
    add_textbox(slide, Inches(8.82), Inches(5.72), Inches(1.5), Inches(0.22), "主动触发", size=10, color=BLUE, bold=True)
    add_textbox(slide, Inches(10.18), Inches(5.72), Inches(1.2), Inches(0.22), "自动触发", size=10, color=GREEN, bold=True)
    add_textbox(slide, Inches(0.5), Inches(1.05), Inches(0.9), Inches(0.28), "长链路", size=10, color=DEEP, bold=True)
    add_textbox(slide, Inches(0.57), Inches(5.28), Inches(0.9), Inches(0.28), "局部任务", size=10, color=DEEP, bold=True)

    add_card(slide, Inches(1.95), Inches(4.15), Inches(2.35), Inches(1.1), "Rule", ["定义边界和禁区"], accent=YELLOW, fill=PALE_YELLOW, tag_text="边界")
    add_card(slide, Inches(2.45), Inches(2.35), Inches(3.1), Inches(1.5), "Skill", ["用户主动发起", "高频工作流稳定复用"], accent=BLUE, fill=LIGHT_BG, tag_text="主动")
    add_card(slide, Inches(6.78), Inches(3.95), Inches(3.0), Inches(1.3), "Hook", ["事件一发生就自动守门", "适合风险拦截"], accent=GREEN, fill=PALE_GREEN, tag_text="自动")
    add_card(slide, Inches(7.15), Inches(1.85), Inches(3.35), Inches(1.75), "Agent", ["跨阶段、跨上下文", "需要更强自主推进能力"], accent=CORAL, fill=PALE_CORAL, tag_text="长链")
    add_band(slide, Inches(2.45), Inches(5.88), Inches(7.7), Inches(0.46), "先想 Skill / Hook，很多问题会比直接做 Agent 更稳。", fill=WHITE, line=SKY, color=INK, bold=True)


def build_p6(slide, section: Section) -> None:
    set_content_header(slide, section.title, "Skill 最适合承接那些会被反复主动发起、流程相对稳定、输出结果有明确形状的任务。")
    quadrants = [
        ("工程类", ["修 bug", "补测试", "做 code review"], BLUE),
        ("文档类", ["写需求", "整理日报", "生成发布说明"], GREEN),
        ("分析类", ["竞品分析", "问题排查", "风险评估"], YELLOW),
        ("协作类", ["会议纪要", "交接说明", "验收总结"], CORAL),
    ]
    positions = [(Inches(0.6), Inches(1.45)), (Inches(6.25), Inches(1.45)), (Inches(0.6), Inches(3.5)), (Inches(6.25), Inches(3.5))]
    for (title, lines, color), (x, y) in zip(quadrants, positions):
        add_card(slide, x, y, Inches(4.95), Inches(1.7), title, lines, accent=color, fill=WHITE, title_size=17, body_size=13)
    add_band(slide, Inches(1.0), Inches(5.55), Inches(9.95), Inches(0.48), "共性特征：主动触发 + 流程稳定 + 输出有明确形状。", fill=LIGHT_BG, line=SKY, color=INK, bold=True)


def build_p7(slide, section: Section) -> None:
    set_content_header(slide, section.title, "Skill 的收益不只是提效，而是把个人经验、团队协作和治理抓手一起沉淀下来。")
    add_circle(slide, Inches(4.78), Inches(2.2), Inches(2.02), "团队\n资产化", fill=BLUE, line=BLUE, size=18)
    orbit_cards = [
        (Inches(1.2), Inches(1.72), "复用", ["减少重复解释", "让同类任务直接复用"], BLUE, LIGHT_BG),
        (Inches(8.1), Inches(1.72), "稳定", ["减少风格漂移", "减少步骤遗漏"], GREEN, PALE_GREEN),
        (Inches(1.2), Inches(4.02), "协作", ["不同角色站在同一流程起点", "输出更容易对齐"], YELLOW, PALE_YELLOW),
        (Inches(8.1), Inches(4.02), "传承", ["把少数人的经验沉淀成团队能力", "让新成员更快接上"], CORAL, PALE_CORAL),
    ]
    for x, y, title, lines, accent, fill in orbit_cards:
        add_card(slide, x, y, Inches(2.9), Inches(1.46), title, lines, accent=accent, fill=fill, title_size=15, body_size=11)
    add_connector(slide, Inches(4.08), Inches(2.52), Inches(4.78), Inches(3.0), color=SOFT_LINE)
    add_connector(slide, Inches(8.1), Inches(2.52), Inches(6.8), Inches(3.0), color=SOFT_LINE)
    add_connector(slide, Inches(4.08), Inches(4.75), Inches(4.78), Inches(4.1), color=SOFT_LINE)
    add_connector(slide, Inches(8.1), Inches(4.75), Inches(6.8), Inches(4.1), color=SOFT_LINE)
    add_band(slide, Inches(2.1), Inches(5.62), Inches(7.8), Inches(0.42), "治理收益：更容易挂接反馈、评估效果并持续维护。", fill=WHITE, line=SKY, color=INK, bold=True)


def build_p8(slide, section: Section) -> None:
    set_content_header(slide, section.title, "做 Skill 之前先看高频、稳定、可验证，避免把一次性任务和噪音流程硬沉淀成资产。")
    criteria_panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.65), Inches(1.42), Inches(4.0), Inches(3.95))
    criteria_panel.adjustments[0] = 0.14
    criteria_panel.fill.solid()
    criteria_panel.fill.fore_color.rgb = LIGHT_BG
    criteria_panel.line.color.rgb = SKY
    criteria_panel.line.width = Pt(1.4)
    add_textbox(slide, Inches(0.9), Inches(1.64), Inches(2.8), Inches(0.3), "值得做的三个判断", size=18, color=BLUE, bold=True)
    criteria = [
        ("01 高频", "是不是会重复出现，值得沉淀"),
        ("02 稳定", "步骤和边界是否相对固定"),
        ("03 可验证", "结果是否容易验收与回看"),
    ]
    for index, (title, desc) in enumerate(criteria):
        y = Inches(2.12 + index * 0.98)
        chip = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.9), y, Inches(0.86), Inches(0.32))
        chip.adjustments[0] = 0.2
        chip.fill.solid()
        chip.fill.fore_color.rgb = BLUE if index == 0 else (GREEN if index == 1 else YELLOW)
        chip.line.fill.background()
        add_textbox(slide, Inches(0.98), y + Inches(0.02), Inches(0.68), Inches(0.2), title, size=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(1.95), y - Inches(0.02), Inches(2.3), Inches(0.42), desc, size=12, color=INK)

    anti = [
        ("反例 1", "一次性任务也硬沉淀成 Skill", CORAL, PALE_CORAL),
        ("反例 2", "边界不清、人人理解不同的流程", YELLOW, PALE_YELLOW),
        ("反例 3", "本该自动执行的守门逻辑却做成 Skill", GREEN, PALE_GREEN),
    ]
    for index, (tag, text, accent, fill) in enumerate(anti):
        add_card(slide, Inches(5.15), Inches(1.58 + index * 1.24), Inches(5.7), Inches(0.96), text, [], accent=accent, fill=fill, tag_text=tag, title_size=13)
    add_band(slide, Inches(5.15), Inches(5.02), Inches(5.7), Inches(0.54), "先做判断，再做制作，能显著降低噪音资产。", fill=WHITE, line=SKY, color=INK, bold=True)


def build_p9(slide, section: Section) -> None:
    add_free_header(slide, section.title, "第一个 Skill 先把触发条件、边界、步骤顺序和输出格式写对，不求一步做大。")
    code_panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(1.25), Inches(5.35), Inches(4.5))
    code_panel.adjustments[0] = 0.14
    code_panel.fill.solid()
    code_panel.fill.fore_color.rgb = DEEP
    code_panel.line.color.rgb = SKY
    code_panel.line.width = Pt(1.3)
    add_textbox(slide, Inches(0.8), Inches(1.48), Inches(1.6), Inches(0.24), "最小 Skill 骨架", size=12, color=SKY, bold=True)
    code_lines = [
        "---",
        "name: fix-bug",
        "description: 在用户描述 bug、报错、失败测试、",
        "  异常回归时使用。先定位根因，再做最小修复，",
        "  并补一个回归验证。",
        "---",
        "",
        "1. 先复述现象、影响范围和复现条件。",
        "2. 搜索相关代码、日志和最近改动，确认根因。",
        "3. 只做最小修复，不顺手重构无关部分。",
        "4. 补一个最贴近故障面的验证。",
        "5. 输出：根因、修改点、验证结果、剩余风险。",
    ]
    add_lines(slide, Inches(0.85), Inches(1.86), Inches(4.8), Inches(3.5), code_lines, size=12, color=WHITE, spacing=2, font_name=MONO)

    steps = [
        ("第一步", "写清什么时候该触发"),
        ("第二步", "写清什么时候不该触发"),
        ("第三步", "写清步骤顺序和输出格式"),
        ("第四步", "写清限制条件，不让它越界"),
    ]
    for index, (label, desc) in enumerate(steps):
        y = Inches(1.48 + index * 0.98)
        color = [BLUE, GREEN, YELLOW, CORAL][index]
        fill = [LIGHT_BG, PALE_GREEN, PALE_YELLOW, PALE_CORAL][index]
        card = add_card(slide, Inches(6.3), y, Inches(4.7), Inches(0.8), desc, [], accent=color, fill=fill, tag_text=label, title_size=12)
        card.line.width = Pt(1.2)
    add_band(slide, Inches(6.3), Inches(5.58), Inches(4.7), Inches(0.46), "第一个版本先求稳，不求全。", fill=WHITE, line=SKY, color=INK, bold=True)


def build_p10(slide, section: Section) -> None:
    set_content_header(slide, section.title, "先定义通用 Skill 本体，再补脚本和参考资料，最后映射到不同宿主入口。")
    layers = [
        ("第一层｜通用 Skill 本体", "通常是一份 <skill-name>/SKILL.md，描述触发条件、步骤顺序、边界和输出。", BLUE, LIGHT_BG, Inches(1.12)),
        ("第二层｜按需补资源", "根据任务复杂度追加 scripts/、references/、assets/，把可执行资源和资料一起带上。", GREEN, PALE_GREEN, Inches(2.58)),
        ("第三层｜宿主适配", "再映射到 Claude、Codex 或其他宿主；入口是适配层，不是 Skill 的唯一真实形态。", YELLOW, PALE_YELLOW, Inches(4.05)),
    ]
    for title, desc, accent, fill, y in layers:
        add_card(slide, Inches(1.25), y, Inches(9.4), Inches(1.05), title, [desc], accent=accent, fill=fill, title_size=16, body_size=12)
    add_textbox(slide, Inches(2.0), Inches(5.45), Inches(1.7), Inches(0.24), "Claude：/skill-name", size=11, color=BLUE, bold=True, align=PP_ALIGN.CENTER, fill_color=WHITE, line_color=BLUE, radius="rounded")
    add_textbox(slide, Inches(5.1), Inches(5.45), Inches(1.7), Inches(0.24), "Codex：$skill-name", size=11, color=GREEN, bold=True, align=PP_ALIGN.CENTER, fill_color=WHITE, line_color=GREEN, radius="rounded")
    add_textbox(slide, Inches(8.2), Inches(5.45), Inches(1.7), Inches(0.24), "其他宿主：适配入口", size=11, color=YELLOW, bold=True, align=PP_ALIGN.CENTER, fill_color=WHITE, line_color=YELLOW, radius="rounded")


def build_p11(slide, section: Section) -> None:
    set_content_header(slide, section.title, "很多反复出错的地方不是语言能力问题，而是项目风格判断没有被固定成流程。")
    add_card(
        slide,
        Inches(0.7),
        Inches(1.52),
        Inches(4.65),
        Inches(3.8),
        "没有 Skill 时",
        ["模块页、示例页、导读页容易混成一种写法", "AI 只看局部段落就开始改", "原本的教学节奏和仓库风格被打乱"],
        accent=CORAL,
        fill=PALE_CORAL,
        title_size=18,
        body_size=13,
    )
    add_card(
        slide,
        Inches(6.65),
        Inches(1.52),
        Inches(4.1),
        Inches(3.8),
        "有 docs-writing-style 之后",
        ["先读仓库规则", "先分页面类型", "再找同类页面参照", "最后才开始写和改"],
        accent=GREEN,
        fill=PALE_GREEN,
        title_size=18,
        body_size=13,
    )
    add_chevron(slide, Inches(5.6), Inches(2.75), Inches(0.62), Inches(0.36), fill=MIST)
    add_band(slide, Inches(1.08), Inches(5.48), Inches(9.75), Inches(0.48), "关键价值：把“写得像这个仓库”从个人感觉变成可执行流程。", fill=WHITE, line=SKY, color=INK, bold=True)


def build_p12(slide, section: Section) -> None:
    set_content_header(slide, section.title, "Skill 既要让不同角色能用，也要让真实使用中的经验能低门槛回流给维护者。")
    layers = [
        ("消费结果", "管理者、外部协作方", BLUE, LIGHT_BG, Inches(1.45)),
        ("使用 Skill", "产品经理、设计师", GREEN, PALE_GREEN, Inches(2.72)),
        ("维护 Skill", "研发、Tech Lead", YELLOW, PALE_YELLOW, Inches(3.99)),
    ]
    for title, desc, accent, fill, y in layers:
        add_card(slide, Inches(1.25), y, Inches(7.9), Inches(0.9), title, [desc], accent=accent, fill=fill, title_size=16, body_size=12)
    add_connector(slide, Inches(5.18), Inches(2.36), Inches(5.18), Inches(2.72), color=SOFT_LINE)
    add_connector(slide, Inches(5.18), Inches(3.63), Inches(5.18), Inches(3.99), color=SOFT_LINE)
    feedback = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(9.45), Inches(1.88), Inches(1.62), Inches(2.95))
    feedback.adjustments[0] = 0.18
    feedback.fill.solid()
    feedback.fill.fore_color.rgb = WHITE
    feedback.line.color.rgb = SKY
    feedback.line.width = Pt(1.4)
    add_textbox(slide, Inches(9.63), Inches(2.04), Inches(1.25), Inches(0.54), "反馈\n通道", size=18, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_lines(slide, Inches(9.6), Inches(2.88), Inches(1.2), Inches(1.38), ["真实问题", "低门槛提交", "经验回写"], size=11, color=INK, spacing=4)
    add_connector(slide, Inches(9.45), Inches(3.32), Inches(9.15), Inches(3.32), color=SOFT_LINE)


def build_p13(slide, section: Section) -> None:
    set_content_header(slide, section.title, "先复用已经打磨好的现成能力，再把真正带有团队约束的那部分做成自建 Skill。")
    tools = [
        ("pdf", ["读取 PDF", "拆分 / 合并", "OCR"], BLUE, LIGHT_BG),
        ("docx", ["创建 Word 文档", "编辑结构", "输出专业文稿"], GREEN, PALE_GREEN),
        ("pptx", ["生成演示稿", "编辑页面", "复用模板"], YELLOW, PALE_YELLOW),
        ("skill-creator", ["从零创建 Skill", "迭代说明", "完善触发规则"], CORAL, PALE_CORAL),
    ]
    for index, (title, lines, accent, fill) in enumerate(tools):
        add_card(slide, Inches(0.7 + index * 2.78), Inches(1.68), Inches(2.35), Inches(2.95), title, lines, accent=accent, fill=fill, title_size=18, body_size=12)
    add_band(slide, Inches(1.2), Inches(5.35), Inches(9.7), Inches(0.58), "判断线：先问“有没有现成能力”，再问“哪些问题真的值得自建”。", fill=WHITE, line=SKY, color=INK, size=13, bold=True)


def build_p14(slide, section: Section) -> None:
    add_free_header(slide, section.title, "三套外部案例最值得借的是结构价值，不是命令数量：它们分别强调角色编排、流程闭环和上下文稳定。")
    columns = [
        ("gstack", "角色编排", ["更像角色化 Skill Pack", "强调多视角显式编排"], BLUE, LIGHT_BG),
        ("superpowers", "流程闭环", ["更像默认工程流水线", "强调设计先行、计划驱动、验证闭环"], GREEN, PALE_GREEN),
        ("GSD", "上下文稳定", ["更像上下文工程层", "强调规格、状态和阶段结论持续稳定"], YELLOW, PALE_YELLOW),
    ]
    for index, (name, tag, lines, accent, fill) in enumerate(columns):
        x = Inches(0.72 + index * 3.73)
        add_card(slide, x, Inches(1.35), Inches(3.18), Inches(3.92), name, lines, accent=accent, fill=fill, title_size=20, body_size=13, tag_text=tag)
    add_band(slide, Inches(0.92), Inches(5.48), Inches(4.8), Inches(0.46), "结论 1：三者不是谁更强，而是解决的问题不同。", fill=WHITE, line=SKY, color=INK, bold=True)
    add_band(slide, Inches(6.08), Inches(5.48), Inches(4.3), Inches(0.46), "结论 2：先借结构，再决定是否引入。", fill=WHITE, line=SKY, color=INK, bold=True)


def build_p15(slide, section: Section) -> None:
    add_free_header(slide, section.title, "团队 30 天落地最稳的做法，是先跑通一个最小闭环，而不是一开始堆很多 Skill。")
    weeks = [
        ("第 1 周", "选 1 个每周至少做两次的高频任务", BLUE),
        ("第 2 周", "选 1 个真正能拦住风险的动作", GREEN),
        ("第 3 周", "建立 1 条低门槛反馈通道", YELLOW),
        ("第 4 周", "复盘频率、稳定性和维护成本", CORAL),
    ]
    add_connector(slide, Inches(1.05), Inches(3.52), Inches(10.72), Inches(3.52), color=SOFT_LINE, width=2)
    for index, (title, desc, accent) in enumerate(weeks):
        x = Inches(1.08 + index * 2.5)
        node_fill = LIGHT_BG if accent == BLUE else (PALE_GREEN if accent == GREEN else (PALE_YELLOW if accent == YELLOW else PALE_CORAL))
        add_circle(slide, x, Inches(2.55), Inches(0.78), str(index + 1), fill=accent, line=accent, size=18)
        add_card(slide, x - Inches(0.2), Inches(3.85), Inches(2.0), Inches(1.24), title, [desc], accent=accent, fill=node_fill, title_size=14, body_size=11)
    add_band(slide, Inches(1.1), Inches(5.48), Inches(9.85), Inches(0.5), "先做 1 个 Skill + 1 个 Hook + 1 条反馈通道，比先造 10 个命令更稳。", fill=WHITE, line=SKY, color=INK, bold=True)


def build_p16(slide, section: Section) -> None:
    set_content_header(slide, section.title, "真正要带走的，不是更多命令，而是更稳定的协作方式。")
    takeaways = [
        ("Takeaway 1", "Skill 解决的是“怎样稳定地做一类事”", BLUE, LIGHT_BG),
        ("Takeaway 2", "先做高频、稳定、可验证的那一小批", GREEN, PALE_GREEN),
        ("Takeaway 3", "先有 Skill 本体，再做宿主适配和组织共享", YELLOW, PALE_YELLOW),
    ]
    for index, (tag, text, accent, fill) in enumerate(takeaways):
        add_card(slide, Inches(0.7), Inches(1.55 + index * 1.2), Inches(5.25), Inches(0.9), text, [], accent=accent, fill=fill, tag_text=tag, title_size=12)
    qa = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(6.4), Inches(1.52), Inches(4.45), Inches(3.92))
    qa.adjustments[0] = 0.16
    qa.fill.solid()
    qa.fill.fore_color.rgb = WHITE
    qa.line.color.rgb = SKY
    qa.line.width = Pt(1.5)
    add_textbox(slide, Inches(6.72), Inches(1.82), Inches(3.8), Inches(0.42), "Q&A", size=24, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(6.72), Inches(2.45), Inches(3.78), Inches(1.06), "你所在团队现在最值得先做成 Skill 的一件事是什么？", size=17, color=DEEP, bold=True, align=PP_ALIGN.CENTER)
    add_lines(slide, Inches(6.88), Inches(3.9), Inches(3.5), Inches(0.96), ["从高频任务找入口", "从风险拦截找 Hook", "从反馈通道开始迭代"], size=12, color=INK, spacing=5)


def build_p17_thanks(slide) -> None:
    return


BUILDERS = {
    1: build_cover,
    2: build_p2,
    3: build_p3,
    4: build_p4,
    5: build_p5,
    6: build_p6,
    7: build_p7,
    8: build_p8,
    9: build_p9,
    10: build_p10,
    11: build_p11,
    12: build_p12,
    13: build_p13,
    14: build_p14,
    15: build_p15,
    16: build_p16,
}


def build_deck(template_path: Path, source_path: Path, output_path: Path) -> Path:
    markdown = source_path.read_text(encoding="utf-8")
    sections = parse_sections(markdown)

    prs = Presentation(str(template_path))
    if prs.slide_width != 12192000 or prs.slide_height != 6858000:
        raise ValueError("Unexpected BOEET slide size.")

    remove_all_slides(prs)
    for number in range(1, 17):
        layout_index = 0 if number == 1 else (5 if number in {5, 9, 14, 15} else 3)
        slide = prs.slides.add_slide(prs.slide_layouts[layout_index])
        BUILDERS[number](slide, sections[number])

    closing = prs.slides.add_slide(prs.slide_layouts[4])
    build_p17_thanks(closing)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return output_path


def main() -> None:
    parser = argparse.ArgumentParser(description="Generate a BOEET-branded prompt-to-skill talk deck.")
    parser.add_argument("--source", type=Path, default=SOURCE_MD)
    parser.add_argument("--template", type=Path, default=DEFAULT_TEMPLATE)
    parser.add_argument("--output", type=Path, default=DEFAULT_OUTPUT)
    args = parser.parse_args()

    if not args.template.exists():
        raise FileNotFoundError(f"Template not found: {args.template}")
    if not args.source.exists():
        raise FileNotFoundError(f"Source markdown not found: {args.source}")

    output = build_deck(args.template, args.source, args.output)
    print(output)


if __name__ == "__main__":
    main()
