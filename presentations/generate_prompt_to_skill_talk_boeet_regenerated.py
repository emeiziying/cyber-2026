#!/usr/bin/env python3
from __future__ import annotations

import argparse
import re
from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR_TYPE
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


REPO_ROOT = Path(__file__).resolve().parents[1]
SOURCE_MD = REPO_ROOT / "presentations" / "prompt-to-skill-talk.md"
OUTPUT_PPTX = REPO_ROOT / "presentations" / "prompt-to-skill-talk-boeet-final.pptx"
TEMPLATE_PPTX = Path.home() / ".agents" / "skills" / "boeet-pptx" / "assets" / "boeet-template-16x9-internal.pptx"


FONT = "微软雅黑"

BLUE = RGBColor(0x00, 0x58, 0x89)
SKY = RGBColor(0x00, 0x9A, 0xD8)
GREEN = RGBColor(0x45, 0xBA, 0x7D)
CORAL = RGBColor(0xFF, 0x8F, 0x82)
YELLOW = RGBColor(0xFF, 0xC7, 0x57)
DARK = RGBColor(0x54, 0x58, 0x5A)
DEEP = RGBColor(0x37, 0x35, 0x45)
INK = RGBColor(0x1F, 0x2D, 0x38)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF4, 0xF8, 0xFB)
MIST = RGBColor(0xE7, 0xF1, 0xF7)
PALE_GREEN = RGBColor(0xE7, 0xF6, 0xEE)
PALE_YELLOW = RGBColor(0xFF, 0xF6, 0xDB)
PALE_CORAL = RGBColor(0xFF, 0xEE, 0xEA)
SOFT_LINE = RGBColor(0xC8, 0xD9, 0xE6)


@dataclass(frozen=True)
class SlideSection:
    slide_no: int
    title: str
    duration: str | None
    goal: str | None
    layout_hint: str | None
    body_points: list[str]
    code_block: str | None
    speaker_notes: str


@dataclass(frozen=True)
class SlideBlueprint:
    layout_name: str
    layout_kind: str
    visual_anchor: str
    summary: str


@dataclass
class SlideModel:
    slide_no: int
    title: str
    core_message: str
    body_points: list[str]
    layout_kind: str
    visual_anchor: str
    speaker_notes: str
    summary: str
    payload: dict[str, object]
    layout_name: str


BLUEPRINTS: dict[int, SlideBlueprint] = {
    1: SlideBlueprint("首页", "cover-progress", "progress-arrow", "面向混合团队的 Skill 方法、案例与落地路径"),
    2: SlideBlueprint("内容页", "pain-chain", "pain-cards-plus-chain", "团队已经会用 AI，但高质量做法仍停留在个人经验、聊天记录和人工兜底里。"),
    3: SlideBlueprint("内容页", "before-after", "two-column-compare", "一次性 Prompt 解决“这次怎么做”；Skill 解决“这类事以后都怎么稳定做”。"),
    4: SlideBlueprint("内容页", "triple-compare", "three-way-contrast", "Prompt 决定这次要什么，Skill 决定怎样稳定做，Rules 决定什么不能错。"),
    5: SlideBlueprint("空白页", "boundary-matrix", "quadrant-matrix", "很多问题不必一上来就做 Agent；先判断它更像边界、主动流程、自动守门，还是长链路任务。"),
    6: SlideBlueprint("内容页", "task-grid", "four-task-categories", "Skill 最适合承接高频、主动发起、流程相对稳定、输出有明确形状的任务。"),
    7: SlideBlueprint("内容页", "value-orbit", "center-hub", "Skill 的收益不只是提效，而是把个人经验、团队协作和治理抓手一起沉淀下来。"),
    8: SlideBlueprint("内容页", "criteria-vs-anti", "decision-panel", "做 Skill 之前先看高频、稳定、可验证，避免把一次性任务和噪音流程硬沉淀成资产。"),
    9: SlideBlueprint("空白页", "code-and-steps", "dark-code-panel", "第一个 Skill 先把触发条件、边界、步骤顺序和输出格式写对，不求一步做大。"),
    10: SlideBlueprint("内容页", "layered-adaptation", "three-layer-stack", "先定义通用 Skill 本体，再补脚本和参考资料，最后映射到不同宿主入口。"),
    11: SlideBlueprint("内容页", "case-contrast", "left-right-case", "很多反复出错的地方不是语言能力问题，而是项目风格判断没有被固定成流程。"),
    12: SlideBlueprint("内容页", "stakeholder-ladder", "three-layer-flow", "Skill 既要让不同角色能用，也要让真实使用中的经验能低门槛回流给维护者。"),
    13: SlideBlueprint("内容页", "tool-cards", "four-tool-cards", "先复用已经打磨好的现成能力，再把真正带有团队约束的部分做成自建 Skill。"),
    14: SlideBlueprint("空白页", "ecosystem-columns", "three-column-map", "三套外部案例最值得借的是结构价值，不是命令数量：它们分别强调角色编排、流程闭环和上下文稳定。"),
    15: SlideBlueprint("空白页", "thirty-day-timeline", "four-week-timeline", "团队 30 天落地最稳的做法，是先跑通一个最小闭环，而不是一开始堆很多 Skill。"),
    16: SlideBlueprint("内容页", "takeaways-qa", "takeaways-plus-qa", "真正要带走的，不是更多命令，而是更稳定的协作方式。"),
}

BODY_SLIDE_COUNT = len(BLUEPRINTS)
TOTAL_SLIDE_COUNT = BODY_SLIDE_COUNT + 1
CLOSING_LAYOUT_NAME = "尾页"


def sanitize_copy(text: str) -> str:
    text = text.strip()
    text = re.sub(r"`([^`]+)`", r"\1", text)
    text = text.replace("“", '"').replace("”", '"').replace("’", "'")
    return text


def split_labeled_point(point: str) -> tuple[str | None, str]:
    cleaned = sanitize_copy(point)
    for sep in ("：", ":"):
        if sep in cleaned:
            left, right = cleaned.split(sep, 1)
            return left.strip(), right.strip()
    return None, cleaned


def collect_sections(markdown: str) -> list[SlideSection]:
    pattern = re.compile(r"^### P(\d+)｜(.+?)\n(.*?)(?=^### P\d+｜|\Z)", re.M | re.S)
    sections: list[SlideSection] = []
    for match in pattern.finditer(markdown):
        slide_no = int(match.group(1))
        title = sanitize_copy(match.group(2))
        block = match.group(3)
        sections.append(parse_section_block(slide_no, title, block))
    expected = list(range(1, 17))
    found = [section.slide_no for section in sections]
    if found != expected:
        raise ValueError(f"Unexpected slide sections: expected {expected}, got {found}")
    return sections


def parse_section_block(slide_no: int, title: str, block: str) -> SlideSection:
    duration = None
    goal = None
    layout_hint = None
    body_points: list[str] = []
    code_block_lines: list[str] = []
    notes_lines: list[str] = []

    lines = block.splitlines()
    i = 0
    while i < len(lines):
        line = lines[i]
        stripped = line.strip()
        if stripped.startswith("- 建议时长："):
            duration = sanitize_copy(stripped.removeprefix("- 建议时长："))
            i += 1
            continue
        if stripped.startswith("- 页面目标："):
            goal = sanitize_copy(stripped.removeprefix("- 页面目标："))
            i += 1
            continue
        if stripped.startswith("- 版式建议："):
            layout_hint = sanitize_copy(stripped.removeprefix("- 版式建议："))
            i += 1
            continue
        if stripped.startswith("- 页内内容："):
            i += 1
            while i < len(lines):
                inner = lines[i]
                inner_stripped = inner.strip()
                if inner.startswith("  - "):
                    body_points.append(sanitize_copy(inner[len("  - ") :]))
                    i += 1
                    continue
                if inner_stripped.startswith("```"):
                    i += 1
                    while i < len(lines) and not lines[i].strip().startswith("```"):
                        code_block_lines.append(lines[i].rstrip("\n"))
                        i += 1
                    if i < len(lines):
                        i += 1
                    continue
                if inner_stripped == "":
                    i += 1
                    continue
                break
            continue
        if stripped.startswith("- 讲稿："):
            i += 1
            while i < len(lines):
                inner = lines[i]
                if inner.startswith(">"):
                    note = inner[1:]
                    if note.startswith(" "):
                        note = note[1:]
                    notes_lines.append(note.rstrip())
                    i += 1
                    continue
                if inner.strip() == "":
                    if notes_lines and notes_lines[-1] != "":
                        notes_lines.append("")
                    i += 1
                    continue
                break
            continue
        i += 1

    notes = "\n".join(notes_lines).strip()
    code_block = "\n".join(code_block_lines).strip() or None
    return SlideSection(
        slide_no=slide_no,
        title=title,
        duration=duration,
        goal=goal,
        layout_hint=layout_hint,
        body_points=body_points,
        code_block=code_block,
        speaker_notes=notes,
    )


def pick_core_message(section: SlideSection) -> str:
    priorities = ("一句话结论", "一句话定义", "结论", "共性特征", "关键价值", "Q&A 提示", "示例", "关键机制")
    for point in section.body_points:
        label, value = split_labeled_point(point)
        if label and any(label.startswith(priority) for priority in priorities):
            return value
    if section.body_points:
        return split_labeled_point(section.body_points[-1])[1]
    return section.title


def body_pairs(section: SlideSection) -> list[tuple[str | None, str]]:
    return [split_labeled_point(point) for point in section.body_points]


def build_payload(section: SlideSection) -> dict[str, object]:
    pairs = body_pairs(section)
    by_label = {label: value for label, value in pairs if label}
    points = [value for _, value in pairs]

    if section.slide_no == 1:
        return {
            "headline": by_label.get("主标题", section.title),
            "subtitle": by_label.get("副标题", BLUEPRINTS[1].summary),
            "audience": by_label.get("受众", ""),
            "conclusion": by_label.get("一句话结论", pick_core_message(section)),
        }
    if section.slide_no == 2:
        cards = []
        for label, value in pairs[:4]:
            cards.append({"tag": label or "", "text": value})
        return {"cards": cards, "conclusion": by_label.get("结论", pick_core_message(section))}
    if section.slide_no == 3:
        before = [value for label, value in pairs if label == "Before"]
        after = [value for label, value in pairs if label == "After"]
        return {"before": before, "after": after, "definition": by_label.get("一句话定义", pick_core_message(section))}
    if section.slide_no == 4:
        columns = pairs[:3]
        return {
            "columns": [{"title": label or f"概念 {index + 1}", "desc": value} for index, (label, value) in enumerate(columns)],
            "bottom": pick_core_message(section),
        }
    if section.slide_no == 5:
        items = [{"title": label or f"对象 {index + 1}", "desc": value} for index, (label, value) in enumerate(pairs[:4])]
        return {"items": items, "bottom": pick_core_message(section)}
    if section.slide_no == 6:
        cards = [{"title": label or f"类别 {index + 1}", "lines": [segment.strip() for segment in value.split("、") if segment.strip()]} for index, (label, value) in enumerate(pairs[:4])]
        return {"cards": cards, "common": by_label.get("共性特征", pick_core_message(section))}
    if section.slide_no == 7:
        values = [{"title": label or f"价值 {index + 1}", "desc": value} for index, (label, value) in enumerate(pairs[:4])]
        return {"values": values, "governance": by_label.get("治理", pick_core_message(section))}
    if section.slide_no == 8:
        criteria = [{"title": label or f"判断 {index + 1}", "desc": value} for index, (label, value) in enumerate(pairs[:3])]
        anti = [{"tag": label or f"反例 {index + 1}", "text": value} for index, (label, value) in enumerate(pairs[3:6])]
        return {"criteria": criteria, "anti": anti}
    if section.slide_no == 9:
        steps = [{"tag": label or f"步骤 {index + 1}", "text": value} for index, (label, value) in enumerate(pairs[:4])]
        return {"steps": steps, "code": section.code_block or "", "bottom": pick_core_message(section)}
    if section.slide_no == 10:
        layers = [{"title": label or f"层 {index + 1}", "desc": value} for index, (label, value) in enumerate(pairs[:3])]
        return {
            "layers": layers,
            "examples": [value for _, value in pairs[3:4]],
            "bottom": pick_core_message(section),
        }
    if section.slide_no == 11:
        return {
            "left_title": "没有 Skill 时",
            "left_lines": [pairs[0][1], pairs[1][1]],
            "right_title": "有 docs-writing-style 之后",
            "right_lines": [pairs[2][1], pairs[3][1]],
            "bottom": pick_core_message(section),
        }
    if section.slide_no == 12:
        layers = [{"title": label or f"层级 {index + 1}", "desc": value} for index, (label, value) in enumerate(pairs[:3])]
        return {"layers": layers, "feedback": by_label.get("关键机制", ""), "bottom": pick_core_message(section)}
    if section.slide_no == 13:
        tools = [{"title": label or f"能力 {index + 1}", "desc": value} for index, (label, value) in enumerate(pairs[:4])]
        return {"tools": tools, "bottom": pick_core_message(section)}
    if section.slide_no == 14:
        columns = [{"title": label or f"方向 {index + 1}", "desc": value} for index, (label, value) in enumerate(pairs[:3])]
        conclusions = [value for _, value in pairs[3:5]]
        return {"columns": columns, "conclusions": conclusions}
    if section.slide_no == 15:
        weeks = [{"title": label or f"阶段 {index + 1}", "desc": value} for index, (label, value) in enumerate(pairs[:4])]
        return {"weeks": weeks, "bottom": pick_core_message(section)}
    if section.slide_no == 16:
        takeaways = [{"tag": label or f"Takeaway {index + 1}", "text": value} for index, (label, value) in enumerate(pairs[:3])]
        return {"takeaways": takeaways, "qa": by_label.get("Q&A 提示", pick_core_message(section))}
    return {"points": points}


def build_models(sections: list[SlideSection]) -> list[SlideModel]:
    models: list[SlideModel] = []
    for section in sections:
        blueprint = BLUEPRINTS[section.slide_no]
        models.append(
            SlideModel(
                slide_no=section.slide_no,
                title=section.title,
                core_message=pick_core_message(section),
                body_points=section.body_points,
                layout_kind=blueprint.layout_kind,
                visual_anchor=blueprint.visual_anchor,
                speaker_notes=section.speaker_notes,
                summary=blueprint.summary,
                payload=build_payload(section),
                layout_name=blueprint.layout_name,
            )
        )
    return models


class DeckRenderer:
    def __init__(self, template_path: Path) -> None:
        self.presentation = Presentation(str(template_path))
        if self.presentation.slide_width != 12192000 or self.presentation.slide_height != 6858000:
            raise ValueError("Unexpected BOEET slide size")
        self.layouts = {layout.name: layout for layout in self.presentation.slide_layouts}
        self._remove_all_slides()

    def _remove_all_slides(self) -> None:
        slide_ids = list(self.presentation.slides._sldIdLst)
        for slide_id in slide_ids:
            self.presentation.part.drop_rel(slide_id.rId)
            self.presentation.slides._sldIdLst.remove(slide_id)

    def render(self, models: list[SlideModel], output_path: Path) -> Path:
        for model in models:
            slide = self.presentation.slides.add_slide(self.layouts[model.layout_name])
            self._dispatch(slide, model)
            self._write_notes(slide, model.speaker_notes)
        self.presentation.slides.add_slide(self.layouts[CLOSING_LAYOUT_NAME])
        output_path.parent.mkdir(parents=True, exist_ok=True)
        self.presentation.save(str(output_path))
        return output_path

    def _dispatch(self, slide, model: SlideModel) -> None:
        dispatch = {
            "cover-progress": self._render_cover,
            "pain-chain": self._render_pain_chain,
            "before-after": self._render_before_after,
            "triple-compare": self._render_triple_compare,
            "boundary-matrix": self._render_boundary_matrix,
            "task-grid": self._render_task_grid,
            "value-orbit": self._render_value_orbit,
            "criteria-vs-anti": self._render_criteria_vs_anti,
            "code-and-steps": self._render_code_and_steps,
            "layered-adaptation": self._render_layered_adaptation,
            "case-contrast": self._render_case_contrast,
            "stakeholder-ladder": self._render_stakeholder_ladder,
            "tool-cards": self._render_tool_cards,
            "ecosystem-columns": self._render_ecosystem_columns,
            "thirty-day-timeline": self._render_timeline,
            "takeaways-qa": self._render_takeaways_qa,
        }
        dispatch[model.layout_kind](slide, model)

    def _write_notes(self, slide, notes: str) -> None:
        notes_frame = slide.notes_slide.notes_text_frame
        notes_frame.clear()
        for index, paragraph_text in enumerate(notes.splitlines()):
            paragraph = notes_frame.paragraphs[0] if index == 0 else notes_frame.add_paragraph()
            run = paragraph.add_run()
            run.text = paragraph_text
            run.font.name = FONT
            run.font.size = Pt(12)
            run.font.color.rgb = DARK

    def _placeholder(self, slide, idx: int):
        for shape in slide.shapes:
            placeholder = getattr(shape, "placeholder_format", None)
            if placeholder and placeholder.idx == idx:
                return shape
        raise KeyError(f"Placeholder {idx} not found")

    def _clear_text(self, shape) -> None:
        if shape.has_text_frame:
            shape.text_frame.clear()

    def _set_placeholder(self, slide, idx: int, text: str, *, size: int, bold: bool = False, color: RGBColor = DARK, align: PP_ALIGN = PP_ALIGN.LEFT) -> None:
        shape = self._placeholder(slide, idx)
        self._clear_text(shape)
        frame = shape.text_frame
        frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        paragraph = frame.paragraphs[0]
        paragraph.alignment = align
        run = paragraph.add_run()
        run.text = text
        self._style(run, size=size, bold=bold, color=color)

    def _content_header(self, slide, title: str, summary: str) -> None:
        self._set_placeholder(slide, 11, title, size=20, bold=True, color=DEEP)
        self._set_placeholder(slide, 13, "", size=11, color=DARK)
        self._add_textbox(slide, Inches(0.48), Inches(0.58), Inches(10.9), Inches(0.24), summary, size=11, color=SKY)

    def _blank_header(self, slide, title: str, summary: str) -> None:
        self._add_textbox(slide, Inches(0.45), Inches(0.18), Inches(11.0), Inches(0.34), title, size=21, color=DEEP, bold=True)
        self._add_textbox(slide, Inches(0.45), Inches(0.56), Inches(10.9), Inches(0.24), summary, size=11, color=SKY)

    def _style(self, run, *, size: int, bold: bool = False, color: RGBColor = DARK, italic: bool = False) -> None:
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color

    def _add_textbox(
        self,
        slide,
        x,
        y,
        w,
        h,
        text: str,
        *,
        size: int = 14,
        bold: bool = False,
        color: RGBColor = DARK,
        align: PP_ALIGN = PP_ALIGN.LEFT,
        fill: RGBColor | None = None,
        line: RGBColor | None = None,
        rounded: bool = False,
        margin: float = 0.08,
    ):
        if rounded:
            box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
            box.adjustments[0] = 0.16
        else:
            box = slide.shapes.add_textbox(x, y, w, h)
        if fill is not None:
            box.fill.solid()
            box.fill.fore_color.rgb = fill
        elif hasattr(box, "fill"):
            box.fill.background()
        if line is not None:
            box.line.color.rgb = line
            box.line.width = Pt(1.2)
        elif hasattr(box, "line"):
            box.line.fill.background()
        frame = box.text_frame
        frame.clear()
        frame.word_wrap = True
        frame.margin_left = Inches(margin)
        frame.margin_right = Inches(margin)
        frame.margin_top = Inches(margin)
        frame.margin_bottom = Inches(margin)
        paragraph = frame.paragraphs[0]
        paragraph.alignment = align
        run = paragraph.add_run()
        run.text = text
        self._style(run, size=size, bold=bold, color=color)
        return box

    def _add_lines(
        self,
        slide,
        x,
        y,
        w,
        h,
        lines: list[str],
        *,
        size: int = 12,
        bold_first: bool = False,
        color: RGBColor = INK,
        spacing: int = 5,
        prefix: str = "",
        fill: RGBColor | None = None,
        line: RGBColor | None = None,
        rounded: bool = False,
        margin: float = 0.05,
    ):
        if rounded:
            box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
            box.adjustments[0] = 0.16
        else:
            box = slide.shapes.add_textbox(x, y, w, h)
        if fill is not None:
            box.fill.solid()
            box.fill.fore_color.rgb = fill
        elif hasattr(box, "fill"):
            box.fill.background()
        if line is not None:
            box.line.color.rgb = line
            box.line.width = Pt(1.2)
        elif hasattr(box, "line"):
            box.line.fill.background()
        frame = box.text_frame
        frame.clear()
        frame.word_wrap = True
        frame.margin_left = Inches(margin)
        frame.margin_right = Inches(margin)
        frame.margin_top = Inches(margin)
        frame.margin_bottom = Inches(margin)
        for index, line_text in enumerate(lines):
            paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            paragraph.alignment = PP_ALIGN.LEFT
            run = paragraph.add_run()
            run.text = f"{prefix}{line_text}" if line_text else ""
            self._style(run, size=size, bold=bold_first and index == 0, color=color)
            paragraph.space_after = Pt(spacing)
        return box

    def _add_card(
        self,
        slide,
        x,
        y,
        w,
        h,
        title: str,
        lines: list[str],
        *,
        accent: RGBColor,
        fill: RGBColor = WHITE,
        title_size: int = 15,
        body_size: int = 11,
        tag: str | None = None,
        title_color: RGBColor | None = None,
    ):
        card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
        card.adjustments[0] = 0.16
        card.fill.solid()
        card.fill.fore_color.rgb = fill
        card.line.color.rgb = accent
        card.line.width = Pt(1.4)
        title_top = y + Inches(0.18)
        if tag:
            tag_shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x + Inches(0.16), y + Inches(0.14), Inches(0.92), Inches(0.28))
            tag_shape.adjustments[0] = 0.18
            tag_shape.fill.solid()
            tag_shape.fill.fore_color.rgb = accent
            tag_shape.line.fill.background()
            self._add_textbox(slide, x + Inches(0.2), y + Inches(0.16), Inches(0.84), Inches(0.2), tag, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
            title_top = y + Inches(0.48)
        self._add_textbox(
            slide,
            x + Inches(0.16),
            title_top,
            w - Inches(0.32),
            Inches(0.34),
            title,
            size=title_size,
            bold=True,
            color=title_color or accent,
        )
        if lines:
            self._add_lines(slide, x + Inches(0.16), title_top + Inches(0.38), w - Inches(0.32), h - Inches(0.56), lines, size=body_size, color=INK, spacing=4)
        return card

    def _add_band(self, slide, x, y, w, h, text: str, *, fill: RGBColor, line: RGBColor, color: RGBColor = INK, size: int = 12, bold: bool = False):
        return self._add_textbox(slide, x, y, w, h, text, size=size, bold=bold, color=color, align=PP_ALIGN.CENTER, fill=fill, line=line, rounded=True, margin=0.06)

    def _add_circle(self, slide, x, y, diameter, text: str, *, fill: RGBColor, line: RGBColor, size: int = 17, color: RGBColor = WHITE):
        circle = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, x, y, diameter, diameter)
        circle.fill.solid()
        circle.fill.fore_color.rgb = fill
        circle.line.color.rgb = line
        circle.line.width = Pt(2)
        self._add_textbox(slide, x + Inches(0.08), y + Inches(0.18), diameter - Inches(0.16), diameter - Inches(0.36), text, size=size, bold=True, color=color, align=PP_ALIGN.CENTER)
        return circle

    def _add_chevron(self, slide, x, y, w, h, *, fill: RGBColor = MIST):
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, x, y, w, h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
        shape.line.fill.background()
        return shape

    def _add_connector(self, slide, x1, y1, x2, y2, *, color: RGBColor = SOFT_LINE, width: float = 1.4):
        connector = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, x1, y1, x2, y2)
        connector.line.color.rgb = color
        connector.line.width = Pt(width)
        return connector

    def _render_cover(self, slide, model: SlideModel) -> None:
        payload = model.payload
        for idx in (10, 13, 23):
            self._set_placeholder(slide, idx, "", size=12)
        self._add_textbox(slide, Inches(0.72), Inches(2.05), Inches(5.0), Inches(1.28), str(payload["headline"]).replace("，", "，\n", 1), size=24, bold=True, color=DEEP)
        self._add_textbox(slide, Inches(0.75), Inches(4.82), Inches(4.55), Inches(0.5), str(payload["subtitle"]), size=16, color=BLUE)
        self._add_textbox(slide, Inches(0.75), Inches(5.28), Inches(4.6), Inches(0.26), str(payload["audience"]), size=13, color=DARK)
        labels = [("Prompt", BLUE), ("Skill", GREEN), ("Team Capability", YELLOW)]
        base_x = Inches(6.42)
        base_y = Inches(2.2)
        for index, (label, accent) in enumerate(labels):
            step = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, base_x + Inches(index * 1.42), base_y, Inches(1.18), Inches(0.56))
            step.adjustments[0] = 0.16
            step.fill.solid()
            step.fill.fore_color.rgb = WHITE
            step.line.color.rgb = accent
            step.line.width = Pt(1.6)
            self._add_textbox(slide, base_x + Inches(index * 1.42) + Inches(0.04), base_y + Inches(0.12), Inches(1.1), Inches(0.2), label, size=11, bold=True, color=accent, align=PP_ALIGN.CENTER)
            if index < len(labels) - 1:
                self._add_chevron(slide, base_x + Inches(index * 1.42) + Inches(1.12), base_y + Inches(0.15), Inches(0.25), Inches(0.24))
        self._add_band(slide, Inches(6.2), Inches(3.25), Inches(4.6), Inches(0.58), str(payload["conclusion"]), fill=LIGHT_BG, line=SKY, bold=True, size=11)
        self._add_textbox(slide, Inches(6.45), Inches(4.36), Inches(3.4), Inches(0.24), "45–60 分钟内部分享", size=11, bold=True, color=SKY)

    def _render_pain_chain(self, slide, model: SlideModel) -> None:
        payload = model.payload
        self._content_header(slide, model.title, model.summary)
        accents = [BLUE, GREEN, YELLOW, CORAL]
        fills = [LIGHT_BG, PALE_GREEN, PALE_YELLOW, PALE_CORAL]
        positions = [(Inches(0.55), Inches(1.38)), (Inches(3.15), Inches(1.38)), (Inches(0.55), Inches(3.18)), (Inches(3.15), Inches(3.18))]
        for index, ((x, y), card) in enumerate(zip(positions, payload["cards"])):
            self._add_card(slide, x, y, Inches(2.35), Inches(1.55), str(card["text"]), [], accent=accents[index], fill=fills[index], tag=str(card["tag"]), title_size=13)
        chain = ["临时 prompt", "临场解释", "结果漂移", "人工兜底"]
        for index, text in enumerate(chain):
            y = Inches(1.58 + index * 0.98)
            self._add_card(slide, Inches(6.15), y, Inches(4.72), Inches(0.72), text, [], accent=[BLUE, SKY, GREEN, CORAL][index], fill=WHITE, title_size=14)
            if index < len(chain) - 1:
                self._add_chevron(slide, Inches(8.25), y + Inches(0.77), Inches(0.42), Inches(0.16))
        self._add_band(slide, Inches(6.15), Inches(5.58), Inches(4.72), Inches(0.58), str(payload["conclusion"]), fill=PALE_YELLOW, line=YELLOW, bold=True, size=11)

    def _render_before_after(self, slide, model: SlideModel) -> None:
        payload = model.payload
        self._content_header(slide, model.title, model.summary)
        self._add_card(slide, Inches(0.72), Inches(1.48), Inches(4.7), Inches(3.55), "Before｜一次性 Prompt", list(payload["before"]), accent=BLUE, fill=LIGHT_BG, title_size=18, body_size=13)
        self._add_card(slide, Inches(6.85), Inches(1.48), Inches(4.1), Inches(3.55), "After｜可复用 Skill", list(payload["after"]), accent=GREEN, fill=PALE_GREEN, title_size=18, body_size=13)
        self._add_chevron(slide, Inches(5.72), Inches(2.78), Inches(0.7), Inches(0.38))
        self._add_band(slide, Inches(1.12), Inches(5.4), Inches(9.88), Inches(0.56), str(payload["definition"]), fill=WHITE, line=SKY, bold=True, size=13, color=DEEP)

    def _render_triple_compare(self, slide, model: SlideModel) -> None:
        payload = model.payload
        self._content_header(slide, model.title, model.summary)
        accents = [BLUE, GREEN, YELLOW]
        fills = [WHITE, PALE_GREEN, PALE_YELLOW]
        for index, column in enumerate(payload["columns"]):
            self._add_card(
                slide,
                Inches(0.62 + index * 3.72),
                Inches(1.56),
                Inches(3.1),
                Inches(3.78),
                str(column["title"]),
                [str(column["desc"])],
                accent=accents[index],
                fill=fills[index],
                title_size=18,
                body_size=12,
            )
        self._add_band(slide, Inches(1.05), Inches(5.56), Inches(9.9), Inches(0.5), str(payload["bottom"]), fill=LIGHT_BG, line=SKY, bold=True)

    def _render_boundary_matrix(self, slide, model: SlideModel) -> None:
        payload = model.payload
        self._blank_header(slide, model.title, model.summary)
        self._add_connector(slide, Inches(1.52), Inches(5.7), Inches(10.9), Inches(5.7), width=1.8)
        self._add_connector(slide, Inches(1.7), Inches(5.92), Inches(1.7), Inches(1.28), width=1.8)
        self._add_chevron(slide, Inches(10.72), Inches(5.57), Inches(0.28), Inches(0.22), fill=SOFT_LINE)
        self._add_chevron(slide, Inches(1.58), Inches(1.15), Inches(0.22), Inches(0.28), fill=SOFT_LINE)
        self._add_textbox(slide, Inches(8.9), Inches(5.82), Inches(1.0), Inches(0.2), "主动触发", size=10, bold=True, color=BLUE)
        self._add_textbox(slide, Inches(10.05), Inches(5.82), Inches(1.0), Inches(0.2), "自动触发", size=10, bold=True, color=GREEN)
        self._add_textbox(slide, Inches(0.55), Inches(1.02), Inches(0.8), Inches(0.2), "长链路", size=10, bold=True, color=DEEP)
        self._add_textbox(slide, Inches(0.62), Inches(5.36), Inches(0.8), Inches(0.2), "局部任务", size=10, bold=True, color=DEEP)
        positions = [
            (Inches(1.9), Inches(4.15), Inches(2.25), Inches(1.05), YELLOW, PALE_YELLOW),
            (Inches(2.52), Inches(2.36), Inches(3.12), Inches(1.48), BLUE, LIGHT_BG),
            (Inches(6.82), Inches(4.0), Inches(3.0), Inches(1.26), GREEN, PALE_GREEN),
            (Inches(7.15), Inches(1.92), Inches(3.42), Inches(1.68), CORAL, PALE_CORAL),
        ]
        for item, position in zip(payload["items"], positions):
            x, y, w, h, accent, fill = position
            self._add_card(slide, x, y, w, h, str(item["title"]), [str(item["desc"])], accent=accent, fill=fill, title_size=16, body_size=11)
        self._add_band(slide, Inches(2.45), Inches(5.98), Inches(7.72), Inches(0.42), str(payload["bottom"]), fill=WHITE, line=SKY, bold=True, size=11)

    def _render_task_grid(self, slide, model: SlideModel) -> None:
        payload = model.payload
        self._content_header(slide, model.title, model.summary)
        accents = [BLUE, GREEN, YELLOW, CORAL]
        fills = [WHITE, PALE_GREEN, PALE_YELLOW, PALE_CORAL]
        positions = [(Inches(0.58), Inches(1.42)), (Inches(6.2), Inches(1.42)), (Inches(0.58), Inches(3.5)), (Inches(6.2), Inches(3.5))]
        for index, ((x, y), card) in enumerate(zip(positions, payload["cards"])):
            self._add_card(slide, x, y, Inches(5.0), Inches(1.72), str(card["title"]), list(card["lines"]), accent=accents[index], fill=fills[index], title_size=16, body_size=13)
        self._add_band(slide, Inches(1.0), Inches(5.56), Inches(9.9), Inches(0.48), str(payload["common"]), fill=LIGHT_BG, line=SKY, bold=True)

    def _render_value_orbit(self, slide, model: SlideModel) -> None:
        payload = model.payload
        self._content_header(slide, model.title, model.summary)
        self._add_circle(slide, Inches(4.8), Inches(2.3), Inches(1.92), "团队\n资产化", fill=BLUE, line=BLUE, size=18)
        positions = [
            (Inches(1.1), Inches(1.72), BLUE, LIGHT_BG),
            (Inches(8.08), Inches(1.72), GREEN, PALE_GREEN),
            (Inches(1.1), Inches(4.02), YELLOW, PALE_YELLOW),
            (Inches(8.08), Inches(4.02), CORAL, PALE_CORAL),
        ]
        for value, (x, y, accent, fill) in zip(payload["values"], positions):
            self._add_card(slide, x, y, Inches(3.0), Inches(1.48), str(value["title"]), [str(value["desc"])], accent=accent, fill=fill, title_size=15, body_size=11)
        self._add_connector(slide, Inches(4.12), Inches(2.54), Inches(4.8), Inches(3.1))
        self._add_connector(slide, Inches(8.08), Inches(2.54), Inches(6.72), Inches(3.1))
        self._add_connector(slide, Inches(4.12), Inches(4.8), Inches(4.8), Inches(4.12))
        self._add_connector(slide, Inches(8.08), Inches(4.8), Inches(6.72), Inches(4.12))
        self._add_band(slide, Inches(2.15), Inches(5.62), Inches(7.85), Inches(0.42), str(payload["governance"]), fill=WHITE, line=SKY, bold=True, size=11)

    def _render_criteria_vs_anti(self, slide, model: SlideModel) -> None:
        payload = model.payload
        self._content_header(slide, model.title, model.summary)
        panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.62), Inches(1.42), Inches(4.0), Inches(3.95))
        panel.adjustments[0] = 0.16
        panel.fill.solid()
        panel.fill.fore_color.rgb = LIGHT_BG
        panel.line.color.rgb = SKY
        panel.line.width = Pt(1.4)
        self._add_textbox(slide, Inches(0.92), Inches(1.66), Inches(2.4), Inches(0.28), "值得做的三个判断", size=18, bold=True, color=BLUE)
        accents = [BLUE, GREEN, YELLOW]
        for index, criterion in enumerate(payload["criteria"]):
            y = Inches(2.12 + index * 0.98)
            chip = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.92), y, Inches(0.92), Inches(0.32))
            chip.adjustments[0] = 0.18
            chip.fill.solid()
            chip.fill.fore_color.rgb = accents[index]
            chip.line.fill.background()
            self._add_textbox(slide, Inches(1.0), y + Inches(0.03), Inches(0.74), Inches(0.18), str(criterion["title"]), size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
            self._add_textbox(slide, Inches(2.02), y - Inches(0.02), Inches(2.3), Inches(0.42), str(criterion["desc"]), size=12, color=INK)
        fills = [PALE_CORAL, PALE_YELLOW, PALE_GREEN]
        accents_right = [CORAL, YELLOW, GREEN]
        for index, anti in enumerate(payload["anti"]):
            self._add_card(slide, Inches(5.15), Inches(1.6 + index * 1.24), Inches(5.75), Inches(0.96), str(anti["text"]), [], accent=accents_right[index], fill=fills[index], tag=str(anti["tag"]), title_size=13)
        self._add_band(slide, Inches(5.15), Inches(5.06), Inches(5.75), Inches(0.5), "先做判断，再做制作，能显著降低噪音资产。", fill=WHITE, line=SKY, bold=True, size=11)

    def _render_code_and_steps(self, slide, model: SlideModel) -> None:
        payload = model.payload
        self._blank_header(slide, model.title, model.summary)
        code_panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.56), Inches(1.24), Inches(5.42), Inches(4.52))
        code_panel.adjustments[0] = 0.14
        code_panel.fill.solid()
        code_panel.fill.fore_color.rgb = DEEP
        code_panel.line.color.rgb = SKY
        code_panel.line.width = Pt(1.2)
        self._add_textbox(slide, Inches(0.84), Inches(1.48), Inches(1.8), Inches(0.2), "最小 Skill 骨架", size=12, bold=True, color=SKY)
        code_lines = [line.rstrip() for line in str(payload["code"]).splitlines()]
        self._add_lines(slide, Inches(0.88), Inches(1.88), Inches(4.92), Inches(3.52), code_lines, size=11, color=WHITE, spacing=2, margin=0.03)
        fills = [LIGHT_BG, PALE_GREEN, PALE_YELLOW, PALE_CORAL]
        accents = [BLUE, GREEN, YELLOW, CORAL]
        for index, step in enumerate(payload["steps"]):
            self._add_card(
                slide,
                Inches(6.3),
                Inches(1.48 + index * 0.98),
                Inches(4.75),
                Inches(0.82),
                str(step["text"]),
                [],
                accent=accents[index],
                fill=fills[index],
                tag=str(step["tag"]),
                title_size=12,
            )
        self._add_band(slide, Inches(6.3), Inches(5.58), Inches(4.75), Inches(0.44), str(payload["bottom"]), fill=WHITE, line=SKY, bold=True, size=11)

    def _render_layered_adaptation(self, slide, model: SlideModel) -> None:
        payload = model.payload
        self._content_header(slide, model.title, model.summary)
        fills = [LIGHT_BG, PALE_GREEN, PALE_YELLOW]
        accents = [BLUE, GREEN, YELLOW]
        for index, layer in enumerate(payload["layers"]):
            self._add_card(slide, Inches(1.18), Inches(1.18 + index * 1.46), Inches(9.5), Inches(1.08), str(layer["title"]), [str(layer["desc"])], accent=accents[index], fill=fills[index], title_size=16, body_size=12)
        self._add_textbox(slide, Inches(2.0), Inches(5.48), Inches(1.82), Inches(0.24), "Claude：/skill-name", size=11, bold=True, color=BLUE, align=PP_ALIGN.CENTER, fill=WHITE, line=BLUE, rounded=True, margin=0.05)
        self._add_textbox(slide, Inches(5.02), Inches(5.48), Inches(1.82), Inches(0.24), "Codex：$skill-name", size=11, bold=True, color=GREEN, align=PP_ALIGN.CENTER, fill=WHITE, line=GREEN, rounded=True, margin=0.05)
        self._add_textbox(slide, Inches(8.08), Inches(5.48), Inches(2.05), Inches(0.24), "其他宿主：适配入口", size=11, bold=True, color=YELLOW, align=PP_ALIGN.CENTER, fill=WHITE, line=YELLOW, rounded=True, margin=0.05)

    def _render_case_contrast(self, slide, model: SlideModel) -> None:
        payload = model.payload
        self._content_header(slide, model.title, model.summary)
        self._add_card(slide, Inches(0.7), Inches(1.52), Inches(4.72), Inches(3.82), str(payload["left_title"]), list(payload["left_lines"]), accent=CORAL, fill=PALE_CORAL, title_size=18, body_size=13)
        self._add_card(slide, Inches(6.55), Inches(1.52), Inches(4.25), Inches(3.82), str(payload["right_title"]), list(payload["right_lines"]), accent=GREEN, fill=PALE_GREEN, title_size=18, body_size=13)
        self._add_chevron(slide, Inches(5.65), Inches(2.78), Inches(0.58), Inches(0.34))
        self._add_band(slide, Inches(1.08), Inches(5.48), Inches(9.78), Inches(0.48), str(payload["bottom"]), fill=WHITE, line=SKY, bold=True)

    def _render_stakeholder_ladder(self, slide, model: SlideModel) -> None:
        payload = model.payload
        self._content_header(slide, model.title, model.summary)
        fills = [LIGHT_BG, PALE_GREEN, PALE_YELLOW]
        accents = [BLUE, GREEN, YELLOW]
        layer_positions = [Inches(1.44), Inches(2.72), Inches(4.0)]
        for layer, y, accent, fill in zip(payload["layers"], layer_positions, accents, fills):
            self._add_card(slide, Inches(1.24), y, Inches(7.92), Inches(0.92), str(layer["title"]), [str(layer["desc"])], accent=accent, fill=fill, title_size=16, body_size=12)
        self._add_connector(slide, Inches(5.2), Inches(2.38), Inches(5.2), Inches(2.72))
        self._add_connector(slide, Inches(5.2), Inches(3.64), Inches(5.2), Inches(4.0))
        self._add_card(slide, Inches(9.42), Inches(1.92), Inches(1.66), Inches(3.0), "反馈\n通道", [str(payload["feedback"]), str(payload["bottom"])], accent=SKY, fill=WHITE, title_size=18, body_size=10, title_color=BLUE)
        self._add_connector(slide, Inches(9.42), Inches(3.36), Inches(9.1), Inches(3.36))

    def _render_tool_cards(self, slide, model: SlideModel) -> None:
        payload = model.payload
        self._content_header(slide, model.title, model.summary)
        fills = [LIGHT_BG, PALE_GREEN, PALE_YELLOW, PALE_CORAL]
        accents = [BLUE, GREEN, YELLOW, CORAL]
        for index, tool in enumerate(payload["tools"]):
            self._add_card(slide, Inches(0.7 + index * 2.78), Inches(1.68), Inches(2.35), Inches(2.96), str(tool["title"]), [str(tool["desc"])], accent=accents[index], fill=fills[index], title_size=18, body_size=12)
        self._add_band(slide, Inches(1.08), Inches(5.36), Inches(9.84), Inches(0.58), str(payload["bottom"]), fill=WHITE, line=SKY, bold=True, size=12)

    def _render_ecosystem_columns(self, slide, model: SlideModel) -> None:
        payload = model.payload
        self._blank_header(slide, model.title, model.summary)
        fills = [LIGHT_BG, PALE_GREEN, PALE_YELLOW]
        accents = [BLUE, GREEN, YELLOW]
        tags = ["角色编排", "流程闭环", "上下文稳定"]
        for index, column in enumerate(payload["columns"]):
            self._add_card(slide, Inches(0.72 + index * 3.73), Inches(1.36), Inches(3.18), Inches(3.94), str(column["title"]), [str(column["desc"])], accent=accents[index], fill=fills[index], title_size=20, body_size=13, tag=tags[index])
        conclusions = list(payload["conclusions"])
        if conclusions:
            self._add_band(slide, Inches(0.92), Inches(5.48), Inches(4.85), Inches(0.46), conclusions[0], fill=WHITE, line=SKY, bold=True, size=11)
        if len(conclusions) > 1:
            self._add_band(slide, Inches(6.0), Inches(5.48), Inches(4.45), Inches(0.46), conclusions[1], fill=WHITE, line=SKY, bold=True, size=11)

    def _render_timeline(self, slide, model: SlideModel) -> None:
        payload = model.payload
        self._blank_header(slide, model.title, model.summary)
        self._add_connector(slide, Inches(1.08), Inches(3.52), Inches(10.78), Inches(3.52), width=2)
        fills = [LIGHT_BG, PALE_GREEN, PALE_YELLOW, PALE_CORAL]
        accents = [BLUE, GREEN, YELLOW, CORAL]
        for index, week in enumerate(payload["weeks"]):
            x = Inches(1.08 + index * 2.48)
            self._add_circle(slide, x, Inches(2.55), Inches(0.78), str(index + 1), fill=accents[index], line=accents[index], size=18)
            self._add_card(slide, x - Inches(0.16), Inches(3.86), Inches(2.05), Inches(1.26), str(week["title"]), [str(week["desc"])], accent=accents[index], fill=fills[index], title_size=14, body_size=11)
        self._add_band(slide, Inches(1.0), Inches(5.48), Inches(10.0), Inches(0.5), str(payload["bottom"]), fill=WHITE, line=SKY, bold=True, size=11)

    def _render_takeaways_qa(self, slide, model: SlideModel) -> None:
        payload = model.payload
        self._content_header(slide, model.title, model.summary)
        fills = [LIGHT_BG, PALE_GREEN, PALE_YELLOW]
        accents = [BLUE, GREEN, YELLOW]
        for index, takeaway in enumerate(payload["takeaways"]):
            self._add_card(slide, Inches(0.7), Inches(1.56 + index * 1.2), Inches(5.3), Inches(0.9), str(takeaway["text"]), [], accent=accents[index], fill=fills[index], tag=str(takeaway["tag"]), title_size=12)
        qa = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(6.4), Inches(1.54), Inches(4.45), Inches(3.94))
        qa.adjustments[0] = 0.16
        qa.fill.solid()
        qa.fill.fore_color.rgb = WHITE
        qa.line.color.rgb = SKY
        qa.line.width = Pt(1.5)
        self._add_textbox(slide, Inches(6.78), Inches(1.84), Inches(3.68), Inches(0.38), "Q&A", size=24, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
        self._add_textbox(slide, Inches(6.7), Inches(2.48), Inches(3.82), Inches(1.02), str(payload["qa"]), size=17, bold=True, color=DEEP, align=PP_ALIGN.CENTER)
        self._add_lines(slide, Inches(6.9), Inches(3.92), Inches(3.48), Inches(0.92), ["从高频任务找入口", "从风险拦截找 Hook", "从反馈通道开始迭代"], size=12, color=INK, spacing=4)


def generate_deck(template_path: Path, source_path: Path, output_path: Path) -> Path:
    markdown = source_path.read_text(encoding="utf-8")
    sections = collect_sections(markdown)
    models = build_models(sections)
    renderer = DeckRenderer(template_path)
    result = renderer.render(models, output_path)
    self_check(result)
    return result


def self_check(output_path: Path) -> None:
    prs = Presentation(str(output_path))
    if len(prs.slides) != TOTAL_SLIDE_COUNT:
        raise ValueError(f"Expected {TOTAL_SLIDE_COUNT} slides, got {len(prs.slides)}")
    if prs.slide_width != 12192000 or prs.slide_height != 6858000:
        raise ValueError("Unexpected slide size after generation")
    slides = list(prs.slides)
    for index, slide in enumerate(slides[:BODY_SLIDE_COUNT], start=1):
        note_text = slide.notes_slide.notes_text_frame.text.strip()
        if not note_text:
            raise ValueError(f"Slide {index} is missing speaker notes")


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Generate a BOEET PPTX deck from prompt-to-skill-talk markdown.")
    parser.add_argument("--source", type=Path, default=SOURCE_MD)
    parser.add_argument("--template", type=Path, default=TEMPLATE_PPTX)
    parser.add_argument("--output", type=Path, default=OUTPUT_PPTX)
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    if not args.source.exists():
        raise FileNotFoundError(f"Source markdown not found: {args.source}")
    if not args.template.exists():
        raise FileNotFoundError(f"BOEET template not found: {args.template}")
    output = generate_deck(args.template, args.source, args.output)
    print(output)


if __name__ == "__main__":
    main()
