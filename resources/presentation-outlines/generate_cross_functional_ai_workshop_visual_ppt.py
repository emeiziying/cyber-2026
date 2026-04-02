from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
FONT = "PingFang SC"

PAPER = RGBColor(247, 243, 236)
INK = RGBColor(24, 38, 53)
MUTED = RGBColor(92, 103, 114)
IVORY = RGBColor(255, 251, 245)
STONE = RGBColor(231, 224, 213)
BRICK = RGBColor(190, 92, 71)
BRICK_DARK = RGBColor(134, 60, 47)
TEAL = RGBColor(45, 112, 113)
TEAL_DARK = RGBColor(27, 86, 87)
GOLD = RGBColor(202, 149, 61)
GREEN = RGBColor(90, 133, 92)
RED = RGBColor(184, 69, 58)
NAVY = RGBColor(20, 34, 49)
WHITE = RGBColor(255, 255, 255)
PALE_GREEN = RGBColor(236, 246, 236)
PALE_RED = RGBColor(251, 237, 233)
PALE_TEAL = RGBColor(234, 245, 245)
PALE_GOLD = RGBColor(249, 243, 224)


def set_bg(slide, color=PAPER):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def set_fill(shape, color, transparency=0.0):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.fill.transparency = transparency


def no_line(shape):
    shape.line.fill.background()


def add_rect(slide, left, top, width, height, color, line_color=None):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height
    )
    set_fill(shape, color)
    if line_color is None:
        no_line(shape)
    else:
        shape.line.color.rgb = line_color
    return shape


def add_round(slide, left, top, width, height, color, line_color=None, transparency=0.0):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    set_fill(shape, color, transparency=transparency)
    if line_color is None:
        no_line(shape)
    else:
        shape.line.color.rgb = line_color
    return shape


def add_oval(slide, left, top, width, height, color, transparency=0.2):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL, left, top, width, height
    )
    set_fill(shape, color, transparency=transparency)
    no_line(shape)
    return shape


def add_text(
    slide,
    left,
    top,
    width,
    height,
    text,
    size=20,
    bold=False,
    color=INK,
    align=PP_ALIGN.LEFT,
    font_name=FONT,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font_name
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return box


def set_box_text(shape, title, subtitle_lines=None, title_color=INK, body_color=MUTED, title_size=18, body_size=12):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Pt(12)
    tf.margin_right = Pt(12)
    tf.margin_top = Pt(10)
    tf.margin_bottom = Pt(10)
    p1 = tf.paragraphs[0]
    r1 = p1.add_run()
    r1.text = title
    r1.font.name = FONT
    r1.font.size = Pt(title_size)
    r1.font.bold = True
    r1.font.color.rgb = title_color
    p1.space_after = Pt(6)
    for line in subtitle_lines or []:
        p = tf.add_paragraph()
        p.bullet = True
        p.space_after = Pt(2)
        r = p.add_run()
        r.text = line
        r.font.name = FONT
        r.font.size = Pt(body_size)
        r.font.color.rgb = body_color


def add_footer(slide, number, color=MUTED):
    add_text(
        slide,
        Inches(12.25),
        Inches(7.0),
        Inches(0.45),
        Inches(0.2),
        str(number),
        size=10,
        color=color,
        align=PP_ALIGN.RIGHT,
    )


def add_header(slide, page, title, minutes, accent=BRICK):
    add_text(slide, Inches(0.72), Inches(0.32), Inches(1.4), Inches(0.24), page, size=11, bold=True, color=accent)
    add_text(slide, Inches(0.72), Inches(0.58), Inches(8.6), Inches(0.48), title, size=25, bold=True, color=INK)
    line = add_rect(slide, Inches(0.72), Inches(1.12), Inches(11.85), Inches(0.04), STONE)
    no_line(line)
    chip = add_round(slide, Inches(11.35), Inches(0.38), Inches(1.15), Inches(0.34), IVORY, line_color=STONE)
    tf = chip.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = minutes
    r.font.name = FONT
    r.font.size = Pt(11)
    r.font.bold = True
    r.font.color.rgb = INK


def add_takeaway(slide, text, color=BRICK_DARK, fill=IVORY):
    shape = add_round(slide, Inches(0.75), Inches(6.55), Inches(11.85), Inches(0.5), fill, line_color=STONE)
    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = "结论："
    r1.font.name = FONT
    r1.font.size = Pt(14)
    r1.font.bold = True
    r1.font.color.rgb = color
    r2 = p.add_run()
    r2.text = text
    r2.font.name = FONT
    r2.font.size = Pt(14)
    r2.font.bold = True
    r2.font.color.rgb = INK


def add_chip(slide, left, top, width, text, fill, text_color=WHITE, line_color=None, size=12):
    shape = add_round(slide, left, top, width, Inches(0.38), fill, line_color=line_color)
    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = text_color
    return shape


def add_section_panel(slide, role, message, accent, side_lines):
    panel = add_round(slide, Inches(0.68), Inches(1.38), Inches(3.0), Inches(4.95), accent, line_color=None)
    add_text(slide, Inches(0.98), Inches(1.72), Inches(1.7), Inches(0.24), "角色页", size=11, bold=True, color=WHITE)
    add_text(slide, Inches(0.98), Inches(2.04), Inches(2.0), Inches(0.6), role, size=30, bold=True, color=WHITE)
    add_text(slide, Inches(0.98), Inches(2.82), Inches(2.35), Inches(1.15), message, size=16, bold=True, color=WHITE)
    y = 4.45
    for line in side_lines:
        add_chip(slide, Inches(0.98), Inches(y), Inches(1.78), line, WHITE, text_color=accent, size=11)
        y += 0.56
    return panel


def create_visual_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Slide 1 Cover
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY)
    add_oval(slide, Inches(8.9), Inches(-0.55), Inches(5.3), Inches(5.3), IVORY, transparency=0.0)
    add_oval(slide, Inches(10.25), Inches(0.48), Inches(2.25), Inches(2.25), BRICK, transparency=0.08)
    add_oval(slide, Inches(9.1), Inches(4.45), Inches(3.3), Inches(3.3), TEAL, transparency=0.12)
    add_rect(slide, Inches(8.5), Inches(0.0), Inches(0.12), Inches(7.5), BRICK)
    add_text(slide, Inches(0.78), Inches(0.95), Inches(3.4), Inches(0.3), "30 分钟综合课", size=12, bold=True, color=RGBColor(230, 218, 199))
    add_text(slide, Inches(0.78), Inches(1.38), Inches(5.4), Inches(1.35), "AI 工作提效入门", size=31, bold=True, color=WHITE)
    add_text(slide, Inches(0.8), Inches(2.55), Inches(5.8), Inches(0.7), "产品、设计、研发都能立刻上手的方法", size=20, color=RGBColor(225, 228, 231))
    add_text(slide, Inches(0.8), Inches(3.18), Inches(5.0), Inches(0.35), "统一认知、统一方法、统一边界", size=14, bold=True, color=RGBColor(231, 208, 188))
    add_chip(slide, Inches(0.82), Inches(4.08), Inches(2.95), "AI 值得在哪些工作里用", BRICK)
    add_chip(slide, Inches(3.95), Inches(4.08), Inches(2.72), "不同角色各自怎么用", TEAL)
    add_chip(slide, Inches(0.82), Inches(4.62), Inches(3.18), "部门内使用 AI 的边界", GOLD, text_color=INK)
    info = add_round(slide, Inches(8.95), Inches(4.22), Inches(3.55), Inches(1.68), IVORY, line_color=RGBColor(226, 214, 198))
    tf = info.text_frame
    tf.clear()
    p1 = tf.paragraphs[0]
    r1 = p1.add_run()
    r1.text = "适用听众"
    r1.font.name = FONT
    r1.font.size = Pt(12)
    r1.font.bold = True
    r1.font.color.rgb = BRICK
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = "产品经理 / 设计师 / 研发 / 少量管理者"
    r2.font.name = FONT
    r2.font.size = Pt(16)
    r2.font.bold = True
    r2.font.color.rgb = INK
    p3 = tf.add_paragraph()
    r3 = p3.add_run()
    r3.text = "统一案例：用户反馈“搜索功能不好用”"
    r3.font.name = FONT
    r3.font.size = Pt(12)
    r3.font.color.rgb = MUTED
    add_footer(slide, 1, color=RGBColor(205, 210, 214))

    # Slide 2
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_oval(slide, Inches(10.7), Inches(-0.25), Inches(2.8), Inches(2.8), BRICK, transparency=0.88)
    add_header(slide, "第 2 页", "为什么现在要学 AI", "2 分钟", accent=BRICK)
    add_round(slide, Inches(0.85), Inches(1.72), Inches(5.25), Inches(3.95), WHITE, line_color=STONE)
    add_text(slide, Inches(1.2), Inches(2.0), Inches(1.3), Inches(0.28), "过去", size=14, bold=True, color=GOLD)
    add_text(slide, Inches(1.18), Inches(2.33), Inches(4.0), Inches(0.62), "AI 更像\n个人效率工具", size=24, bold=True, color=INK)
    add_text(slide, Inches(1.2), Inches(3.45), Inches(4.4), Inches(1.3), "• 自己摸索，个人用个人的\n• 重点是提速写初稿\n• 信息仍然要手工在角色间传递", size=16, color=MUTED)
    add_round(slide, Inches(7.02), Inches(1.72), Inches(5.25), Inches(3.95), WHITE, line_color=STONE)
    add_text(slide, Inches(7.37), Inches(2.0), Inches(1.3), Inches(0.28), "现在", size=14, bold=True, color=TEAL)
    add_text(slide, Inches(7.35), Inches(2.33), Inches(4.3), Inches(0.62), "AI 正在变成\n协作基础设施", size=24, bold=True, color=INK)
    add_text(slide, Inches(7.38), Inches(3.45), Inches(4.5), Inches(1.3), "• 更快整理上下文\n• 更快形成可交接材料\n• 更快推进产品、设计、研发协同", size=16, color=MUTED)
    arrow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, Inches(6.12), Inches(3.0), Inches(0.55), Inches(0.95))
    set_fill(arrow, STONE)
    no_line(arrow)
    add_takeaway(slide, "AI 改变的不是某一个岗位，而是整个部门处理信息和交接任务的方式。")
    add_footer(slide, 2)

    # Slide 3
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_header(slide, "第 3 页", "AI 适合做什么，不适合直接放行什么", "3 分钟", accent=RED)
    left_bg = add_round(slide, Inches(0.82), Inches(1.72), Inches(5.8), Inches(4.7), PALE_GREEN, line_color=None)
    right_bg = add_round(slide, Inches(6.72), Inches(1.72), Inches(5.8), Inches(4.7), PALE_RED, line_color=None)
    add_text(slide, Inches(1.15), Inches(2.02), Inches(2.0), Inches(0.35), "适合交给 AI 的工作", size=18, bold=True, color=GREEN)
    add_text(slide, Inches(1.15), Inches(2.55), Inches(4.5), Inches(2.35), "整理\n归纳\n初稿\n对比\n清单\n复盘", size=22, bold=True, color=INK)
    add_text(slide, Inches(7.05), Inches(2.02), Inches(3.1), Inches(0.35), "不适合直接放行的工作", size=18, bold=True, color=RED)
    add_text(slide, Inches(7.05), Inches(2.55), Inches(4.75), Inches(2.65), "事实判断\n业务拍板\n敏感信息处理\n最终设计 / 需求 / 代码决策", size=21, bold=True, color=INK)
    divider = add_round(slide, Inches(5.82), Inches(2.4), Inches(1.68), Inches(1.68), IVORY, line_color=STONE)
    tf = divider.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "先划\n边界"
    r.font.name = FONT
    r.font.size = Pt(18)
    r.font.bold = True
    r.font.color.rgb = BRICK_DARK
    add_takeaway(slide, "AI 可以帮你准备，但不能替你承担责任。", color=RED, fill=RGBColor(255, 247, 244))
    add_footer(slide, 3)

    # Slide 4
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_oval(slide, Inches(-0.8), Inches(5.35), Inches(2.9), Inches(2.9), TEAL, transparency=0.9)
    add_header(slide, "第 4 页", "所有角色都通用的 5 段式输入法", "2 分钟", accent=TEAL)
    steps = [
        ("01", "目标", "你想得到什么", BRICK),
        ("02", "上下文", "这件事发生在什么场景", TEAL),
        ("03", "约束", "哪些边界不能碰", GOLD),
        ("04", "输出格式", "希望它怎么呈现", GREEN),
        ("05", "参考样例", "有没有现成材料", BRICK_DARK),
    ]
    line = add_rect(slide, Inches(1.2), Inches(2.45), Inches(10.85), Inches(0.06), STONE)
    no_line(line)
    x_positions = [1.05, 3.35, 5.65, 7.95, 10.25]
    for idx, (num, label, desc, color) in enumerate(steps):
        circle = add_oval(slide, Inches(x_positions[idx]), Inches(1.75), Inches(1.0), Inches(1.0), color, transparency=0.0)
        tf = circle.text_frame
        tf.clear()
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = num
        r.font.name = FONT
        r.font.size = Pt(18)
        r.font.bold = True
        r.font.color.rgb = WHITE
        add_text(slide, Inches(x_positions[idx] - 0.15), Inches(2.9), Inches(1.3), Inches(0.24), label, size=14, bold=True, color=INK, align=PP_ALIGN.CENTER)
        add_text(slide, Inches(x_positions[idx] - 0.35), Inches(3.25), Inches(1.7), Inches(0.62), desc, size=12, color=MUTED, align=PP_ALIGN.CENTER)
    sample = add_round(slide, Inches(0.95), Inches(4.75), Inches(11.4), Inches(1.15), IVORY, line_color=STONE)
    tf = sample.text_frame
    tf.clear()
    p1 = tf.paragraphs[0]
    r1 = p1.add_run()
    r1.text = "示例输入"
    r1.font.name = FONT
    r1.font.size = Pt(12)
    r1.font.bold = True
    r1.font.color.rgb = TEAL_DARK
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = "请把最近两周的搜索相关用户反馈整理成 3 类问题，并给出优先级建议。不要直接给解决方案，只输出问题分类表。"
    r2.font.name = FONT
    r2.font.size = Pt(14)
    r2.font.color.rgb = INK
    add_takeaway(slide, "输入越完整，输出越稳定。", color=TEAL_DARK)
    add_footer(slide, 4)

    # Slide 5
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_header(slide, "第 5 页", "所有角色都通用的 4 步验收法", "2 分钟", accent=GOLD)
    left_panel = add_round(slide, Inches(0.82), Inches(1.65), Inches(3.5), Inches(4.9), WHITE, line_color=STONE)
    add_text(slide, Inches(1.16), Inches(2.02), Inches(2.6), Inches(0.32), "4 步验收法", size=16, bold=True, color=GOLD)
    add_text(slide, Inches(1.14), Inches(2.42), Inches(2.8), Inches(1.0), "不要只看\n“写得像不像”", size=25, bold=True, color=INK)
    add_text(slide, Inches(1.16), Inches(3.75), Inches(2.8), Inches(1.4), "要看：\n• 有没有依据\n• 有没有遗漏边界\n• 能不能直接放行", size=15, color=MUTED)
    items = [
        ("1", "有无依据", "结论有没有来源", BRICK),
        ("2", "边界完整", "有没有漏掉限制条件", TEAL),
        ("3", "格式正确", "是否真的符合要求", GOLD),
        ("4", "需要复核", "能不能直接放行", GREEN),
    ]
    coords = [(4.65, 1.8), (8.55, 1.8), (4.65, 4.15), (8.55, 4.15)]
    for (num, title, desc, color), (left, top) in zip(items, coords):
        card = add_round(slide, Inches(left), Inches(top), Inches(3.0), Inches(1.75), WHITE, line_color=STONE)
        badge = add_oval(slide, Inches(left + 0.18), Inches(top + 0.2), Inches(0.75), Inches(0.75), color, transparency=0.0)
        tf = badge.text_frame
        tf.clear()
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = num
        r.font.name = FONT
        r.font.size = Pt(15)
        r.font.bold = True
        r.font.color.rgb = WHITE
        add_text(slide, Inches(left + 1.05), Inches(top + 0.28), Inches(1.7), Inches(0.24), title, size=15, bold=True, color=INK)
        add_text(slide, Inches(left + 1.05), Inches(top + 0.72), Inches(1.7), Inches(0.4), desc, size=11, color=MUTED)
    quote = add_round(slide, Inches(0.9), Inches(6.0), Inches(11.6), Inches(0.48), PALE_GOLD, line_color=STONE)
    tf = quote.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "案例追问：如果 AI 说“搜索不好用是因为排序问题”，要继续问：依据是什么？有没有用户原话或数据？"
    r.font.name = FONT
    r.font.size = Pt(13)
    r.font.color.rgb = INK
    add_takeaway(slide, "验收不是研发专属动作，所有角色都需要。", color=GOLD)
    add_footer(slide, 5)

    # Slide 6
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_header(slide, "第 6 页", "产品经理怎么用 AI", "3 分钟", accent=BRICK)
    add_section_panel(
        slide,
        "产品经理",
        "最适合让 AI 做\n结构化整理\n和边界补全",
        BRICK,
        ["需求澄清", "反馈归类", "验收草稿"],
    )
    cards = [
        ("需求澄清", ["把模糊想法整理成需求提纲"]),
        ("反馈归类", ["把工单和访谈归成问题类型"]),
        ("竞品框架", ["先出对比结构，再人工补事实"]),
        ("验收标准草稿", ["先列边界和失败场景"]),
        ("会议纪要行动项", ["自动拆成负责人和截止时间"]),
    ]
    positions = [(4.05, 1.7), (7.15, 1.7), (10.25, 1.7), (4.05, 3.95), (7.15, 3.95)]
    widths = [2.8, 2.8, 2.2, 2.8, 5.35]
    for (title, lines), (left, top), width in zip(cards, positions, widths):
        shape = add_round(slide, Inches(left), Inches(top), Inches(width), Inches(1.7), WHITE, line_color=STONE)
        set_box_text(shape, title, lines, title_color=INK, title_size=16, body_size=12)
    case = add_round(slide, Inches(10.25), Inches(3.95), Inches(2.2), Inches(1.7), PALE_RED, line_color=None)
    set_box_text(case, "统一案例", ["把“搜索不好用”整理成：查不到 / 查不准 / 筛选难理解。"], title_color=BRICK_DARK, body_color=INK, title_size=15, body_size=11)
    add_takeaway(slide, "产品不应把需求判断直接外包给 AI。", color=BRICK_DARK, fill=RGBColor(255, 245, 241))
    add_footer(slide, 6)

    # Slide 7
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_header(slide, "第 7 页", "设计师怎么用 AI", "3 分钟", accent=TEAL)
    add_section_panel(
        slide,
        "设计师",
        "最适合把 AI 当\n评审助手\n和表达加速器",
        TEAL,
        ["方向发散", "交互说明", "评审清单"],
    )
    ribbon = add_round(slide, Inches(4.05), Inches(1.74), Inches(8.35), Inches(1.12), PALE_TEAL, line_color=None)
    add_chip(slide, Inches(4.3), Inches(2.03), Inches(2.1), "方向发散", TEAL)
    add_chip(slide, Inches(7.0), Inches(2.03), Inches(2.1), "交互说明", BRICK, size=12)
    add_chip(slide, Inches(9.7), Inches(2.03), Inches(2.1), "评审清单", GOLD, text_color=INK, size=12)
    chevron1 = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, Inches(6.45), Inches(2.04), Inches(0.38), Inches(0.32))
    set_fill(chevron1, STONE)
    no_line(chevron1)
    chevron2 = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, Inches(9.15), Inches(2.04), Inches(0.38), Inches(0.32))
    set_fill(chevron2, STONE)
    no_line(chevron2)
    info_cards = [
        ("界面文案草稿", ["空态、错误态、按钮文案"]),
        ("设计评审清单", ["一致性、可用性、信息层级"]),
        ("系统命名 / 状态枚举", ["先出结构，再人工定稿"]),
        ("统一案例", ["检查搜索页是否缺少空结果、无权限、无筛选条件等状态说明。"]),
    ]
    positions = [(4.05, 3.35), (7.22, 3.35), (4.05, 5.03), (8.12, 5.03)]
    widths = [2.85, 2.85, 3.65, 4.28]
    fills = [WHITE, WHITE, WHITE, PALE_TEAL]
    titles = [INK, INK, INK, TEAL_DARK]
    for (title, lines), (left, top), width, fill, title_color in zip(info_cards, positions, widths, fills, titles):
        shape = add_round(slide, Inches(left), Inches(top), Inches(width), Inches(1.35), fill, line_color=STONE if fill == WHITE else None)
        set_box_text(shape, title, lines, title_color=title_color, body_color=MUTED if fill == WHITE else INK, title_size=15, body_size=11)
    add_takeaway(slide, "设计不能把审美判断直接外包给 AI。", color=TEAL_DARK, fill=RGBColor(241, 250, 250))
    add_footer(slide, 7)

    # Slide 8
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY)
    add_oval(slide, Inches(9.9), Inches(-0.8), Inches(4.0), Inches(4.0), BRICK, transparency=0.83)
    add_oval(slide, Inches(-1.0), Inches(4.8), Inches(3.0), Inches(3.0), TEAL, transparency=0.82)
    add_text(slide, Inches(0.72), Inches(0.3), Inches(1.6), Inches(0.24), "第 8 页", size=11, bold=True, color=RGBColor(224, 206, 182))
    add_text(slide, Inches(0.72), Inches(0.58), Inches(8.0), Inches(0.48), "研发怎么用 AI", size=25, bold=True, color=WHITE)
    add_chip(slide, Inches(11.35), Inches(0.38), Inches(1.1), "2 分钟", IVORY, text_color=INK)
    add_text(slide, Inches(0.9), Inches(1.52), Inches(3.1), Inches(0.8), "研发的 AI 使用更深，\n但本质仍然是给上下文 + 做验证。", size=20, bold=True, color=WHITE)
    quad = [
        ("读代码", "快速定位相关逻辑", BRICK),
        ("补测试", "为改动补回归测试", GREEN),
        ("做 Review", "提前发现边界遗漏", TEAL),
        ("批量重构", "处理重复性修改", GOLD),
    ]
    coords = [(4.25, 1.7), (8.45, 1.7), (4.25, 4.1), (8.45, 4.1)]
    for (title, body, color), (left, top) in zip(quad, coords):
        shape = add_round(slide, Inches(left), Inches(top), Inches(3.45), Inches(1.75), RGBColor(28, 46, 66), line_color=RGBColor(57, 72, 89))
        badge = add_rect(slide, Inches(left), Inches(top), Inches(0.12), Inches(1.75), color)
        no_line(badge)
        set_box_text(shape, title, [body], title_color=WHITE, body_color=RGBColor(209, 216, 221), title_size=17, body_size=12)
    add_round(slide, Inches(0.88), Inches(4.45), Inches(2.6), Inches(1.08), RGBColor(28, 46, 66), line_color=RGBColor(57, 72, 89))
    add_text(slide, Inches(1.12), Inches(4.75), Inches(2.1), Inches(0.3), "围绕搜索功能", size=12, bold=True, color=RGBColor(224, 206, 182))
    add_text(slide, Inches(1.1), Inches(5.05), Inches(2.2), Inches(0.6), "AI 可以帮助研发更快找到筛选、排序、关键词匹配相关代码位置。", size=11, color=RGBColor(209, 216, 221))
    add_footer(slide, 8, color=RGBColor(205, 210, 214))

    # Slide 9
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_header(slide, "第 9 页", "一个需求如何跨角色流转", "4 分钟", accent=BRICK)
    add_text(slide, Inches(0.9), Inches(1.55), Inches(4.1), Inches(0.24), "统一案例：用户反馈“搜索功能不好用”", size=14, bold=True, color=BRICK)
    columns = [
        ("产品", BRICK, ["整理反馈", "归类问题", "补充优先级"]),
        ("设计", TEAL, ["补交互状态", "整理说明", "准备评审清单"]),
        ("研发", GOLD, ["定位逻辑", "修改实现", "补测试"]),
        ("验收", GREEN, ["统一用 4 步", "检查是否可用"]),
    ]
    lefts = [0.9, 3.93, 6.96, 9.99]
    for idx, (title, color, lines) in enumerate(columns):
        card = add_round(slide, Inches(lefts[idx]), Inches(2.1), Inches(2.45), Inches(3.8), WHITE, line_color=STONE)
        top_band = add_rect(slide, Inches(lefts[idx]), Inches(2.1), Inches(2.45), Inches(0.18), color)
        no_line(top_band)
        add_text(slide, Inches(lefts[idx] + 0.2), Inches(2.45), Inches(2.0), Inches(0.32), title, size=18, bold=True, color=color)
        y = 3.0
        for line in lines:
            bubble = add_round(slide, Inches(lefts[idx] + 0.18), Inches(y), Inches(2.05), Inches(0.52), IVORY, line_color=STONE)
            tf = bubble.text_frame
            tf.clear()
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = line
            r.font.name = FONT
            r.font.size = Pt(12)
            r.font.bold = True
            r.font.color.rgb = INK
            y += 0.72
        if idx < 3:
            arrow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, Inches(lefts[idx] + 2.48), Inches(3.55), Inches(0.35), Inches(0.55))
            set_fill(arrow, STONE)
            no_line(arrow)
    role_bar = add_round(slide, Inches(0.95), Inches(6.0), Inches(11.45), Inches(0.38), PALE_GOLD, line_color=None)
    tf = role_bar.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "AI 在每一段的角色：产品阶段整理材料 / 设计阶段补表达和检查清单 / 研发阶段读上下文和验证实现"
    r.font.name = FONT
    r.font.size = Pt(12)
    r.font.bold = True
    r.font.color.rgb = INK
    add_takeaway(slide, "AI 的最大价值是减少角色之间“重新解释一遍”的成本。", color=BRICK_DARK)
    add_footer(slide, 9)

    # Slide 10
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, RGBColor(42, 31, 30))
    add_oval(slide, Inches(9.5), Inches(-0.6), Inches(4.6), Inches(4.6), RED, transparency=0.84)
    add_text(slide, Inches(0.74), Inches(0.34), Inches(1.7), Inches(0.24), "第 10 页", size=11, bold=True, color=RGBColor(236, 218, 213))
    add_text(slide, Inches(0.74), Inches(0.6), Inches(8.0), Inches(0.48), "部门内使用 AI 的红线", size=25, bold=True, color=WHITE)
    add_chip(slide, Inches(11.35), Inches(0.38), Inches(1.1), "2 分钟", RGBColor(77, 57, 53), text_color=WHITE)
    add_text(slide, Inches(0.94), Inches(1.56), Inches(2.8), Inches(1.15), "5 条红线", size=38, bold=True, color=RGBColor(244, 209, 201))
    redlines = [
        "敏感数据直接输入公共工具",
        "外部账号、客户信息、生产环境配置直接暴露",
        "未经审阅的内容直接对外发布",
        "带权限的自动化操作直接放行",
        "把 AI 的判断当作正式结论",
    ]
    y = 2.15
    for text in redlines:
        box = add_round(slide, Inches(4.0), Inches(y), Inches(8.0), Inches(0.68), RGBColor(68, 50, 47), line_color=RGBColor(100, 76, 72))
        tf = box.text_frame
        tf.clear()
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = text
        r.font.name = FONT
        r.font.size = Pt(14)
        r.font.bold = True
        r.font.color.rgb = WHITE
        y += 0.82
    advice = add_round(slide, Inches(0.96), Inches(5.25), Inches(2.5), Inches(1.05), RGBColor(68, 50, 47), line_color=RGBColor(100, 76, 72))
    tf = advice.text_frame
    tf.clear()
    p1 = tf.paragraphs[0]
    r1 = p1.add_run()
    r1.text = "建议做法"
    r1.font.name = FONT
    r1.font.size = Pt(12)
    r1.font.bold = True
    r1.font.color.rgb = RGBColor(244, 209, 201)
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = "边界不清楚时，先按高风险处理。"
    r2.font.name = FONT
    r2.font.size = Pt(13)
    r2.font.bold = True
    r2.font.color.rgb = WHITE
    add_footer(slide, 10, color=RGBColor(210, 205, 205))

    # Slide 11
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_oval(slide, Inches(-0.6), Inches(-0.7), Inches(3.0), Inches(3.0), GOLD, transparency=0.9)
    add_header(slide, "第 11 页", "部门落地建议", "2 分钟", accent=GREEN)
    add_text(slide, Inches(0.95), Inches(1.68), Inches(4.2), Inches(0.42), "先试点，再沉淀，再推广", size=24, bold=True, color=INK)
    roadmap = [
        ("01", "全员通识", "统一认知和边界", GOLD),
        ("02", "角色试点", "各挑 1-2 个高频场景", BRICK),
        ("03", "统一模板", "沉淀输入、验收、复盘模板", TEAL),
        ("04", "周期复盘", "持续更新规范", GREEN),
    ]
    x = 0.95
    for idx, (num, title, desc, color) in enumerate(roadmap):
        card = add_round(slide, Inches(x), Inches(2.45), Inches(2.75), Inches(2.1), WHITE, line_color=STONE)
        badge = add_oval(slide, Inches(x + 0.18), Inches(2.65), Inches(0.68), Inches(0.68), color, transparency=0.0)
        tf = badge.text_frame
        tf.clear()
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = num
        r.font.name = FONT
        r.font.size = Pt(12)
        r.font.bold = True
        r.font.color.rgb = WHITE if color != GOLD else INK
        add_text(slide, Inches(x + 1.0), Inches(2.76), Inches(1.5), Inches(0.25), title, size=16, bold=True, color=INK)
        add_text(slide, Inches(x + 0.2), Inches(3.5), Inches(2.3), Inches(0.6), desc, size=12, color=MUTED)
        if idx < 3:
            arrow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, Inches(x + 2.58), Inches(3.15), Inches(0.42), Inches(0.55))
            set_fill(arrow, STONE)
            no_line(arrow)
        x += 2.92
    add_chip(slide, Inches(1.0), Inches(5.35), Inches(1.6), "会议纪要", GREEN, text_color=WHITE)
    add_chip(slide, Inches(2.85), Inches(5.35), Inches(2.15), "用户反馈归类", BRICK)
    add_chip(slide, Inches(5.3), Inches(5.35), Inches(2.1), "设计说明整理", TEAL)
    add_chip(slide, Inches(7.7), Inches(5.35), Inches(1.6), "测试补充", GOLD, text_color=INK)
    add_takeaway(slide, "不要一开始追求全员全面铺开，先把高频低风险场景跑通。", color=GREEN, fill=RGBColor(244, 249, 244))
    add_footer(slide, 11)

    # Slide 12
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY)
    add_oval(slide, Inches(-0.85), Inches(-0.65), Inches(3.8), Inches(3.8), TEAL, transparency=0.85)
    add_oval(slide, Inches(10.2), Inches(4.2), Inches(3.4), Inches(3.4), BRICK, transparency=0.86)
    add_text(slide, Inches(0.78), Inches(0.95), Inches(2.6), Inches(0.3), "最后记住 3 句话", size=12, bold=True, color=RGBColor(224, 206, 182))
    statements = [
        "AI 最适合做准备工作、结构化工作、重复工作",
        "AI 用得越多，越需要边界、验收和复盘",
        "一个部门里最值得优化的是角色之间的交接方式",
    ]
    y = 1.75
    for idx, text in enumerate(statements, start=1):
        add_text(slide, Inches(0.95), Inches(y), Inches(0.55), Inches(0.38), f"{idx:02d}", size=15, bold=True, color=RGBColor(224, 206, 182))
        add_text(slide, Inches(1.55), Inches(y - 0.02), Inches(8.9), Inches(0.52), text, size=22, bold=True, color=WHITE)
        line = add_rect(slide, Inches(1.55), Inches(y + 0.55), Inches(8.6), Inches(0.03), RGBColor(63, 79, 96))
        no_line(line)
        y += 1.45
    q = add_round(slide, Inches(10.0), Inches(5.55), Inches(2.35), Inches(0.85), IVORY, line_color=None)
    tf = q.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "Q&A"
    r.font.name = FONT
    r.font.size = Pt(20)
    r.font.bold = True
    r.font.color.rgb = INK
    add_footer(slide, 12, color=RGBColor(205, 210, 214))

    return prs


if __name__ == "__main__":
    path = Path(__file__).with_name("30min-cross-functional-ai-workshop-visual.pptx")
    prs = create_visual_presentation()
    prs.save(path)
    print(path)
