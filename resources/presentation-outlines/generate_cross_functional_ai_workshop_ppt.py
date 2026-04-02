from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

BG = RGBColor(248, 245, 238)
NAVY = RGBColor(19, 35, 52)
SLATE = RGBColor(88, 100, 112)
ACCENT = RGBColor(198, 91, 66)
ACCENT_DARK = RGBColor(146, 58, 39)
TEAL = RGBColor(49, 110, 112)
GREEN = RGBColor(79, 128, 85)
RED = RGBColor(181, 67, 55)
GOLD = RGBColor(202, 145, 51)
CARD = RGBColor(255, 255, 255)
CARD_ALT = RGBColor(241, 236, 226)
LINE = RGBColor(219, 211, 199)

FONT = "PingFang SC"


def set_background(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text="", size=20, bold=False,
                color=NAVY, align=PP_ALIGN.LEFT, font_name=FONT):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Pt(6)
    frame.margin_right = Pt(6)
    frame.margin_top = Pt(4)
    frame.margin_bottom = Pt(4)
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_paragraphs(frame, lines, size=18, color=NAVY, level=0, space_after=5):
    for idx, line in enumerate(lines):
        paragraph = frame.paragraphs[0] if idx == 0 and not frame.paragraphs[0].text else frame.add_paragraph()
        paragraph.level = level
        paragraph.space_after = Pt(space_after)
        run = paragraph.add_run()
        run.text = line
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.color.rgb = color


def add_title(slide, section, title, minutes, accent=ACCENT):
    add_textbox(slide, Inches(0.6), Inches(0.28), Inches(2.0), Inches(0.3),
                section, size=13, bold=True, color=accent)
    add_textbox(slide, Inches(0.6), Inches(0.52), Inches(8.7), Inches(0.6),
                title, size=26, bold=True, color=NAVY)
    pill = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(11.45), Inches(0.38), Inches(1.05), Inches(0.34)
    )
    pill.fill.solid()
    pill.fill.fore_color.rgb = CARD_ALT
    pill.line.color.rgb = LINE
    pill.text_frame.text = minutes
    paragraph = pill.text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.runs[0]
    run.font.name = FONT
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = NAVY
    pill.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.6), Inches(1.12), Inches(12.1), Inches(0.03)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = LINE
    line.line.fill.background()


def add_footer(slide, page_number):
    add_textbox(
        slide, Inches(12.35), Inches(7.0), Inches(0.35), Inches(0.2),
        str(page_number), size=10, color=SLATE, align=PP_ALIGN.RIGHT
    )


def add_card(slide, left, top, width, height, title, bullets=None,
             fill_color=CARD, title_color=NAVY, bullet_color=SLATE, accent_color=ACCENT):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = LINE
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Pt(10)
    tf.margin_right = Pt(10)
    tf.margin_top = Pt(10)
    tf.margin_bottom = Pt(8)
    p1 = tf.paragraphs[0]
    r1 = p1.add_run()
    r1.text = title
    r1.font.name = FONT
    r1.font.size = Pt(15)
    r1.font.bold = True
    r1.font.color.rgb = title_color
    p1.space_after = Pt(6)
    if bullets:
        for item in bullets:
            p = tf.add_paragraph()
            p.bullet = True
            p.space_after = Pt(2)
            r = p.add_run()
            r.text = item
            r.font.name = FONT
            r.font.size = Pt(12)
            r.font.color.rgb = bullet_color
    tag = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        left, top, Inches(0.08), height
    )
    tag.fill.solid()
    tag.fill.fore_color.rgb = accent_color
    tag.line.fill.background()
    return shape


def add_takeaway(slide, text, color=ACCENT_DARK):
    box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.7), Inches(6.55), Inches(11.95), Inches(0.45)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = CARD_ALT
    box.line.color.rgb = LINE
    tf = box.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "结论：" + text
    r.font.name = FONT
    r.font.size = Pt(14)
    r.font.bold = True
    r.font.color.rgb = color


def create_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Slide 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.0))
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    band.line.fill.background()
    add_textbox(slide, Inches(0.7), Inches(1.3), Inches(8.8), Inches(0.8),
                "AI 工作提效入门", size=30, bold=True, color=NAVY)
    add_textbox(slide, Inches(0.7), Inches(2.0), Inches(9.7), Inches(0.5),
                "产品、设计、研发都能立刻上手的方法", size=20, color=SLATE)
    add_textbox(slide, Inches(0.7), Inches(2.65), Inches(8.5), Inches(0.42),
                "30 分钟建立统一认知、统一方法、统一边界", size=14, bold=True, color=ACCENT)
    pills = [
        ("AI 值得在哪些工作里用", Inches(0.75), GREEN),
        ("不同角色各自怎么用", Inches(4.35), TEAL),
        ("部门内使用 AI 的边界", Inches(7.95), ACCENT),
    ]
    for text, left, color in pills:
        pill = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            left, Inches(3.55), Inches(3.05), Inches(0.62)
        )
        pill.fill.solid()
        pill.fill.fore_color.rgb = CARD
        pill.line.color.rgb = color
        tf = pill.text_frame
        tf.clear()
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = text
        r.font.name = FONT
        r.font.size = Pt(13)
        r.font.bold = True
        r.font.color.rgb = color
    add_textbox(slide, Inches(0.8), Inches(5.15), Inches(5.4), Inches(0.35),
                "适用听众", size=13, bold=True, color=ACCENT)
    add_textbox(slide, Inches(0.8), Inches(5.45), Inches(5.6), Inches(0.6),
                "产品经理 / 设计师 / 研发 / 少量管理者", size=18, color=NAVY)
    add_textbox(slide, Inches(7.6), Inches(5.15), Inches(4.5), Inches(0.35),
                "统一案例", size=13, bold=True, color=ACCENT)
    add_textbox(slide, Inches(7.6), Inches(5.45), Inches(4.8), Inches(0.7),
                "用户反馈：搜索功能不好用", size=18, color=NAVY)
    add_footer(slide, 1)

    # Slide 2
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "第 2 页", "为什么现在要学 AI", "2 分钟")
    add_card(
        slide, Inches(0.8), Inches(1.55), Inches(5.7), Inches(3.85),
        "过去：AI 更像个人效率工具",
        ["自己摸索，个人用个人的", "重点是提速写初稿", "信息仍然要手工在角色间传递"],
        fill_color=CARD, accent_color=GOLD
    )
    add_card(
        slide, Inches(6.8), Inches(1.55), Inches(5.7), Inches(3.85),
        "现在：AI 正在变成协作基础设施",
        ["更快整理上下文", "更快形成可交接材料", "更快推进产品、设计、研发协同"],
        fill_color=CARD, accent_color=TEAL
    )
    add_takeaway(slide, "AI 改变的不是某一个岗位，而是整个部门处理信息和交接任务的方式。")
    add_footer(slide, 2)

    # Slide 3
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "第 3 页", "AI 适合做什么，不适合直接放行什么", "3 分钟")
    add_card(
        slide, Inches(0.8), Inches(1.55), Inches(5.85), Inches(4.7),
        "适合交给 AI 的工作",
        ["整理", "归纳", "初稿", "对比", "清单", "复盘"],
        fill_color=RGBColor(244, 250, 244), title_color=GREEN, bullet_color=NAVY, accent_color=GREEN
    )
    add_card(
        slide, Inches(6.72), Inches(1.55), Inches(5.85), Inches(4.7),
        "不适合直接放行的工作",
        ["事实判断", "业务拍板", "敏感信息处理", "最终设计决策", "最终需求决策", "最终代码决策"],
        fill_color=RGBColor(252, 242, 239), title_color=RED, bullet_color=NAVY, accent_color=RED
    )
    add_takeaway(slide, "AI 可以帮你准备，但不能替你承担责任。", color=RED)
    add_footer(slide, 3)

    # Slide 4
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "第 4 页", "所有角色都通用的 5 段式输入法", "2 分钟")
    labels = [
        ("01 目标", "你想得到什么"),
        ("02 上下文", "事情发生在什么场景"),
        ("03 约束", "哪些边界不能碰"),
        ("04 输出格式", "希望它怎么呈现"),
        ("05 参考样例", "有没有现成材料可参考"),
    ]
    x_positions = [Inches(0.8), Inches(3.25), Inches(5.7), Inches(8.15), Inches(10.6)]
    for idx, (title, desc) in enumerate(labels):
        add_card(
            slide, x_positions[idx], Inches(1.8), Inches(2.1), Inches(2.7),
            title, [desc], fill_color=CARD, accent_color=[ACCENT, TEAL, GOLD, GREEN, ACCENT_DARK][idx]
        )
    sample = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.82), Inches(5.0), Inches(11.75), Inches(1.1)
    )
    sample.fill.solid()
    sample.fill.fore_color.rgb = CARD_ALT
    sample.line.color.rgb = LINE
    tf = sample.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "示例：请把最近两周的搜索相关用户反馈整理成 3 类问题，并给出优先级建议。不要直接给解决方案，只输出问题分类表。"
    r.font.name = FONT
    r.font.size = Pt(14)
    r.font.color.rgb = NAVY
    add_takeaway(slide, "输入越完整，输出越稳定。")
    add_footer(slide, 4)

    # Slide 5
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "第 5 页", "所有角色都通用的 4 步验收法", "2 分钟")
    checks = [
        ("1. 有无依据", "结论有没有来源"),
        ("2. 边界完整", "有没有漏掉限制条件"),
        ("3. 格式正确", "是否真的符合要求"),
        ("4. 需要复核", "能不能直接放行"),
    ]
    for idx, item in enumerate(checks):
        add_card(
            slide, Inches(0.9 + idx * 3.05), Inches(2.0), Inches(2.75), Inches(2.5),
            item[0], [item[1]], fill_color=CARD, accent_color=[GOLD, TEAL, ACCENT, GREEN][idx]
        )
    quote = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(1.0), Inches(5.1), Inches(11.3), Inches(0.9)
    )
    quote.fill.solid()
    quote.fill.fore_color.rgb = RGBColor(245, 248, 252)
    quote.line.color.rgb = LINE
    tf = quote.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "案例追问：如果 AI 说“搜索不好用是因为排序问题”，要继续问：依据是什么？有没有用户原话或数据？"
    r.font.name = FONT
    r.font.size = Pt(15)
    r.font.color.rgb = NAVY
    add_takeaway(slide, "不要只看它写得像不像，要看它能不能用、敢不敢用。")
    add_footer(slide, 5)

    # Slide 6
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "第 6 页", "产品经理怎么用 AI", "3 分钟")
    pm_cards = [
        ("需求澄清", "把模糊想法整理成需求提纲"),
        ("反馈归类", "把工单和访谈归成问题类型"),
        ("竞品框架", "先出对比结构，再人工补事实"),
        ("验收草稿", "先列边界和失败场景"),
        ("会议纪要", "自动拆行动项和负责人"),
    ]
    coords = [(0.8, 1.7), (3.35, 1.7), (5.9, 1.7), (8.45, 1.7), (2.1, 4.2)]
    widths = [2.3, 2.3, 2.3, 2.3, 4.4]
    for idx, (title, desc) in enumerate(pm_cards):
        left, top = coords[idx]
        width = widths[idx]
        add_card(
            slide, Inches(left), Inches(top), Inches(width), Inches(1.75),
            title, [desc], fill_color=CARD, accent_color=ACCENT
        )
    add_card(
        slide, Inches(7.0), Inches(4.2), Inches(5.3), Inches(1.75),
        "统一案例带入",
        ["把“搜索不好用”的反馈整理成：查不到 / 查不准 / 筛选难理解 三类。"],
        fill_color=RGBColor(250, 239, 235), accent_color=ACCENT_DARK
    )
    add_takeaway(slide, "产品最适合让 AI 做结构化整理和边界补全，不适合把需求判断直接外包。")
    add_footer(slide, 6)

    # Slide 7
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "第 7 页", "设计师怎么用 AI", "3 分钟")
    design_cards = [
        ("方向发散", "快速列多个界面方向"),
        ("文案草稿", "空态、错误态、按钮文案"),
        ("评审清单", "一致性、可用性、信息层级"),
        ("交互说明", "把页面说明写得更利于开发理解"),
        ("系统命名", "组件状态和命名结构草稿"),
    ]
    coords = [(0.8, 1.7), (3.35, 1.7), (5.9, 1.7), (8.45, 1.7), (2.1, 4.2)]
    widths = [2.3, 2.3, 2.3, 2.3, 4.4]
    for idx, (title, desc) in enumerate(design_cards):
        left, top = coords[idx]
        width = widths[idx]
        add_card(
            slide, Inches(left), Inches(top), Inches(width), Inches(1.75),
            title, [desc], fill_color=CARD, accent_color=TEAL
        )
    add_card(
        slide, Inches(7.0), Inches(4.2), Inches(5.3), Inches(1.75),
        "统一案例带入",
        ["检查搜索页是否缺少空结果、无权限、无筛选条件等状态说明。"],
        fill_color=RGBColor(236, 247, 247), accent_color=TEAL
    )
    add_takeaway(slide, "设计最适合把 AI 当评审助手和表达加速器，而不是审美决策者。", color=TEAL)
    add_footer(slide, 7)

    # Slide 8
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "第 8 页", "研发怎么用 AI", "2 分钟")
    eng_cards = [
        ("读代码", "快速定位相关逻辑"),
        ("补测试", "为改动补回归测试"),
        ("做 Review", "提前发现边界遗漏"),
        ("批量重构", "处理重复性修改"),
    ]
    coords = [(0.9, 1.9), (6.7, 1.9), (0.9, 4.0), (6.7, 4.0)]
    accents = [ACCENT, GREEN, TEAL, GOLD]
    for idx, ((title, desc), (left, top)) in enumerate(zip(eng_cards, coords)):
        add_card(
            slide, Inches(left), Inches(top), Inches(5.5), Inches(1.55),
            title, [desc], fill_color=CARD, accent_color=accents[idx]
        )
    add_takeaway(slide, "研发的 AI 使用更深，但本质仍然是给上下文 + 做验证。")
    add_footer(slide, 8)

    # Slide 9
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "第 9 页", "一个需求如何跨角色流转", "4 分钟")
    roles = [
        ("产品", "整理反馈\n归类问题\n补充优先级", ACCENT),
        ("设计", "补交互状态\n整理说明\n准备评审清单", TEAL),
        ("研发", "定位逻辑\n修改实现\n补测试", GOLD),
        ("验收", "统一用 4 步\n检查是否可用", GREEN),
    ]
    lefts = [0.8, 3.75, 6.7, 9.65]
    for idx, (title, body, color) in enumerate(roles):
        card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(lefts[idx]), Inches(2.1), Inches(2.25), Inches(2.9)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = CARD
        card.line.color.rgb = LINE
        tf = card.text_frame
        tf.clear()
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.alignment = PP_ALIGN.CENTER
        r1 = p1.add_run()
        r1.text = title
        r1.font.name = FONT
        r1.font.size = Pt(18)
        r1.font.bold = True
        r1.font.color.rgb = color
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = body
        r2.font.name = FONT
        r2.font.size = Pt(14)
        r2.font.color.rgb = NAVY
        if idx < len(roles) - 1:
            arrow = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.CHEVRON,
                Inches(lefts[idx] + 2.3), Inches(3.05), Inches(0.45), Inches(0.7)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = LINE
            arrow.line.fill.background()
    add_textbox(slide, Inches(0.9), Inches(1.55), Inches(4.4), Inches(0.3),
                "统一案例：用户反馈“搜索功能不好用”", size=14, bold=True, color=ACCENT)
    add_takeaway(slide, "AI 的最大价值不只是单点提效，而是减少角色之间重新解释一遍的成本。", color=ACCENT_DARK)
    add_footer(slide, 9)

    # Slide 10
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "第 10 页", "部门内使用 AI 的红线", "2 分钟")
    add_card(
        slide, Inches(0.85), Inches(1.7), Inches(7.2), Inches(4.4),
        "这些情况默认不能直接放行",
        [
            "敏感数据直接输入公共工具",
            "外部账号、客户信息、生产环境配置直接暴露",
            "未经审阅的内容直接对外发布",
            "带权限的自动化操作直接放行",
            "把 AI 的判断当作正式结论",
        ],
        fill_color=RGBColor(252, 242, 239), title_color=RED, bullet_color=NAVY, accent_color=RED
    )
    add_card(
        slide, Inches(8.35), Inches(2.0), Inches(4.0), Inches(2.7),
        "建议做法",
        ["高风险事项默认人工确认", "边界不清楚时，先按高风险处理", "先问“能不能给 AI”，再问“能不能更快”"],
        fill_color=CARD, accent_color=ACCENT_DARK
    )
    add_takeaway(slide, "边界先于效率。", color=RED)
    add_footer(slide, 10)

    # Slide 11
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "第 11 页", "部门落地建议", "2 分钟")
    roadmap = [
        ("1. 全员通识", "统一认知和边界"),
        ("2. 角色试点", "各挑 1-2 个高频场景"),
        ("3. 统一模板", "沉淀输入、验收、复盘模板"),
        ("4. 周期复盘", "持续更新规范"),
    ]
    x = 0.8
    widths = [2.75, 2.75, 2.95, 2.8]
    accents = [GOLD, ACCENT, TEAL, GREEN]
    for idx, ((title, desc), width) in enumerate(zip(roadmap, widths)):
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.CHEVRON,
            Inches(x), Inches(2.1), Inches(width), Inches(1.55)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = accents[idx]
        shape.line.fill.background()
        tf = shape.text_frame
        tf.clear()
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.alignment = PP_ALIGN.CENTER
        r1 = p1.add_run()
        r1.text = title
        r1.font.name = FONT
        r1.font.size = Pt(16)
        r1.font.bold = True
        r1.font.color.rgb = CARD
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = desc
        r2.font.name = FONT
        r2.font.size = Pt(11)
        r2.font.color.rgb = CARD
        x += width - 0.15
    add_card(
        slide, Inches(1.05), Inches(4.45), Inches(11.1), Inches(1.2),
        "可以先试点的低风险场景",
        ["会议纪要", "用户反馈归类", "设计说明整理", "测试补充"],
        fill_color=CARD, accent_color=GREEN
    )
    add_takeaway(slide, "不要一开始追求全员全面铺开，先把高频低风险场景跑通。", color=GREEN)
    add_footer(slide, 11)

    # Slide 12
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "第 12 页", "总结与问答", "2 分钟")
    statements = [
        ("AI 最适合做", "准备工作 / 结构化工作 / 重复工作", ACCENT),
        ("最终判断仍归", "人", TEAL),
        ("部门里最值得优化的", "是角色之间的交接方式", GOLD),
    ]
    top = 1.8
    for title, value, color in statements:
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(0.9), Inches(top), Inches(11.5), Inches(1.0)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = CARD
        shape.line.color.rgb = LINE
        tf = shape.text_frame
        tf.clear()
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        r1 = p.add_run()
        r1.text = title + "："
        r1.font.name = FONT
        r1.font.size = Pt(18)
        r1.font.bold = True
        r1.font.color.rgb = color
        r2 = p.add_run()
        r2.text = value
        r2.font.name = FONT
        r2.font.size = Pt(18)
        r2.font.bold = True
        r2.font.color.rgb = NAVY
        top += 1.2
    add_card(
        slide, Inches(8.7), Inches(5.7), Inches(3.4), Inches(0.9),
        "Q&A 引导", ["你最想先试哪个场景？", "你最担心的风险是什么？"],
        fill_color=CARD_ALT, accent_color=ACCENT_DARK
    )
    add_footer(slide, 12)

    return prs


if __name__ == "__main__":
    output = Path(__file__).with_name("30min-cross-functional-ai-workshop.pptx")
    presentation = create_presentation()
    presentation.save(output)
    print(output)
